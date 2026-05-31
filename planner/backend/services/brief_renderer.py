"""
PowerPoint brief generator — squadron template token substitution.

A template is a normal .pptx file containing `{{token.path}}` placeholders
anywhere a text run can live: slide titles, body text, tables, grouped
shapes. The frontend resolves each token to a value from the mission store
and posts back the value map; this service does the actual substitution
and returns rendered bytes.

Design choices:
  - Stateless. No session lookup. The /api/brief/render endpoint takes
    the template + a flat {token: value} dict and returns the filled .pptx.
    This keeps the brief generator independent of mission-data shape
    changes — the frontend is the only place that knows how to resolve
    a token to a value.
  - Tokens are pure string substitution. We don't try to type-check or
    eval expressions. `{{flight[0].callsign}}` is just a string the
    frontend resolved to `'BENGAL11'` before sending it.
  - Tokens missing from the values map are LEFT AS-IS (`{{token}}` in the
    output) so the user can spot what didn't get filled.
  - Run-level formatting is preserved when a token sits inside a single
    run. PowerPoint sometimes splits `{{ flight.callsign }}` across multiple
    runs (especially after manual editing); when that happens we fall back
    to a paragraph-level substitution that loses intra-paragraph formatting
    on that one paragraph but keeps the substitution working.
"""
import io
import math
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Token paths must start with a letter and may contain dots, brackets,
# digits, underscores. e.g. {{mission.theater}}, {{flight[0].callsign}},
# {{weather.wind.surface_kt}}.
TOKEN_RE = re.compile(r"\{\{\s*([a-zA-Z][a-zA-Z0-9_.\[\]]*)\s*\}\}")


def scan_template(template_bytes: bytes) -> List[str]:
    """Return a sorted list of unique token paths found in the template.

    Raises ValueError on malformed .pptx (corrupted ZIP, bad XML, etc.).
    """
    try:
        prs = Presentation(io.BytesIO(template_bytes))
    except Exception as e:
        raise ValueError(f"Could not open .pptx: {e}") from e

    tokens: set = set()
    for slide in prs.slides:
        _collect_tokens_from_shapes(slide.shapes, tokens)
        # Speaker notes can also contain tokens — useful for handouts.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                for run in paragraph.runs:
                    for m in TOKEN_RE.finditer(run.text):
                        tokens.add(m.group(1))
    return sorted(tokens)


def render_template(template_bytes: bytes, values: Dict[str, str]) -> bytes:
    """Substitute every recognised token in the template with values[token].

    Tokens missing from `values` are left as literal `{{token}}` so the
    user can spot them in the rendered output and decide whether to
    add them to the value map.

    Returns the filled .pptx as raw bytes.
    """
    try:
        prs = Presentation(io.BytesIO(template_bytes))
    except Exception as e:
        raise ValueError(f"Could not open .pptx: {e}") from e

    for slide in prs.slides:
        _substitute_in_shapes(slide.shapes, values)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            _substitute_in_text_frame(slide.notes_slide.notes_text_frame, values)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# --- internals ---------------------------------------------------------------

def _collect_tokens_from_shapes(shapes, tokens: set) -> None:
    for shape in shapes:
        if shape.has_text_frame:
            _collect_from_text_frame(shape.text_frame, tokens)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _collect_from_text_frame(cell.text_frame, tokens)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_tokens_from_shapes(shape.shapes, tokens)


def _collect_from_text_frame(text_frame, tokens: set) -> None:
    # Use the joined-paragraph approach so split-run tokens are still found.
    for paragraph in text_frame.paragraphs:
        full = "".join(run.text for run in paragraph.runs)
        for m in TOKEN_RE.finditer(full):
            tokens.add(m.group(1))


def _substitute_in_shapes(shapes, values: Dict[str, str]) -> None:
    for shape in shapes:
        if shape.has_text_frame:
            _substitute_in_text_frame(shape.text_frame, values)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _substitute_in_text_frame(cell.text_frame, values)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _substitute_in_shapes(shape.shapes, values)


def _substitute_in_text_frame(text_frame, values: Dict[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs
        if not runs:
            continue

        # Fast path: token contained entirely within a single run.
        # Preserves font / size / bold across the rest of the paragraph.
        all_in_run = True
        modified = False
        for run in runs:
            if TOKEN_RE.search(run.text):
                modified = True
                # If the token spans into the next run, this regex won't match
                # the whole thing — flag for slow-path fallback.
                # We only treat it as fast-path-clean if every {{ has a }} on
                # the same run.
                if run.text.count("{{") != run.text.count("}}"):
                    all_in_run = False
                    break

        if modified and all_in_run:
            for run in runs:
                if TOKEN_RE.search(run.text):
                    run.text = TOKEN_RE.sub(
                        lambda m: str(values.get(m.group(1), m.group(0))),
                        run.text,
                    )
            continue

        # Slow path: token spans runs OR no token at all. Either way, do
        # paragraph-level substitution and write back into run 0.
        full = "".join(run.text for run in runs)
        if not TOKEN_RE.search(full):
            continue
        substituted = TOKEN_RE.sub(
            lambda m: str(values.get(m.group(1), m.group(0))),
            full,
        )
        runs[0].text = substituted
        for run in runs[1:]:
            run.text = ""


# --- Format conversion (PPTX → PDF / PNG / JPG) ----------------------------
#
# Uses LibreOffice headless (`soffice --convert-to`) which is free, runs
# server-side, and produces high-fidelity output. The Dockerfile installs
# `libreoffice-impress + libreoffice-core` for production.
#
# For local dev: install LibreOffice (https://www.libreoffice.org/download/)
# and ensure `soffice` is on PATH. Without it, `convert_pptx` raises
# `LibreOfficeNotFoundError` which the API surfaces as a 503 with a
# helpful message — the rest of the brief workflow still works (.pptx
# download path is unaffected).

ConvertFormat = Literal["pdf", "png", "jpg", "pptx"]


class LibreOfficeNotFoundError(RuntimeError):
    """Raised when soffice/libreoffice can't be located on the system."""


def _find_soffice() -> Optional[str]:
    """Locate the LibreOffice CLI binary across Linux/macOS/Windows."""
    # Prefer explicit env override if set (handy for non-standard installs).
    override = os.environ.get("LIBREOFFICE_PATH")
    if override and os.path.isfile(override):
        return override

    # PATH lookup — covers Linux/macOS standard installs and Docker.
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found

    # Windows default install paths — `where soffice` doesn't find it
    # because the installer doesn't add it to PATH by default.
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for path in candidates:
        if os.path.isfile(path):
            return path
    return None


def is_conversion_available() -> bool:
    """Cheap check the API can use to decide whether to advertise the option."""
    return _find_soffice() is not None


def convert_pptx(pptx_bytes: bytes, target: ConvertFormat) -> Tuple[bytes, str]:
    """Convert a rendered .pptx to PDF, PNG (zip of slides), or JPG (zip).

    Returns (output_bytes, mime_type). `target='pptx'` is a no-op pass-through
    so callers can route any format through one function.

    PNG/JPG outputs are bundled into a ZIP because LibreOffice emits one
    image file per slide and the API needs a single response body.

    Raises:
      LibreOfficeNotFoundError — soffice not installed
      RuntimeError — conversion subprocess failed (with stderr in message)
      TimeoutError — conversion exceeded 60s
    """
    if target == "pptx":
        return pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    # Validate target up-front so a bad value gives a clean ValueError
    # regardless of whether LibreOffice is installed.
    if target not in ("pdf", "png", "jpg"):
        raise ValueError(f"Unsupported target format: {target}")

    soffice = _find_soffice()
    if soffice is None:
        raise LibreOfficeNotFoundError(
            "LibreOffice is not installed on this server. "
            "Install it (https://www.libreoffice.org/download/) and ensure "
            "`soffice` is on PATH, or set LIBREOFFICE_PATH to its location. "
            "PowerPoint download (.pptx) still works without LibreOffice."
        )

    # All non-pptx targets go through a single pptx→pdf soffice call.
    # PDF is returned directly; PNG/JPG are rasterized per-page from the PDF.
    # Running soffice once (instead of twice) avoids Windows file-lock contention.
    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        src = tmpdir / "brief.pptx"
        src.write_bytes(pptx_bytes)

        # Per-call user profile so we don't contend with any other soffice
        # process on the host. Path.as_uri() is cross-platform — produces
        # `file:///C:/...` on Windows, `file:///home/...` on Linux.
        user_profile = tmpdir / "soffice_profile"
        cmd = [
            soffice,
            "--headless",
            "--norestore",
            "--nologo",
            "--nofirststartwizard",
            f"-env:UserInstallation={user_profile.as_uri()}",
            "--convert-to", "pdf",
            "--outdir", str(tmpdir),
            str(src),
        ]
        try:
            result = subprocess.run(
                cmd, capture_output=True, timeout=60, check=False,
            )
        except subprocess.TimeoutExpired as e:
            raise TimeoutError("LibreOffice conversion timed out (60s)") from e

        if result.returncode != 0:
            stderr = result.stderr.decode("utf-8", errors="replace")[:500]
            raise RuntimeError(f"LibreOffice conversion failed: {stderr}")

        pdf_path = tmpdir / "brief.pdf"
        if not pdf_path.exists():
            raise RuntimeError("Expected brief.pdf not produced by LibreOffice")
        pdf_bytes = pdf_path.read_bytes()

        if target == "pdf":
            return pdf_bytes, "application/pdf"

        # PNG/JPG — rasterize each PDF page and zip the images.
        return _pdf_to_image_zip(pdf_bytes, target)


def _pdf_to_image_zip(
    pdf_bytes: bytes, target: Literal["png", "jpg"],
) -> Tuple[bytes, str]:
    """Rasterize a PDF to per-page images and bundle them as a ZIP."""
    images: List[bytes] = _rasterize_pdf(pdf_bytes, target, dpi=150)
    if not images:
        raise RuntimeError("No images produced from PDF rasterization")

    buf = io.BytesIO()
    ext = "png" if target == "png" else "jpg"
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, img_bytes in enumerate(images, start=1):
            zf.writestr(f"slide_{i:02d}.{ext}", img_bytes)
    return buf.getvalue(), "application/zip"


def _rasterize_pdf(pdf_bytes: bytes, target: Literal["png", "jpg"], dpi: int = 150) -> List[bytes]:
    """Rasterize a PDF to per-page images. Tries pypdfium2 first."""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return _rasterize_pdf_via_ghostscript(pdf_bytes, target, dpi)

    pdf = pdfium.PdfDocument(pdf_bytes)
    out: List[bytes] = []
    scale = dpi / 72  # PDF default is 72dpi
    for page in pdf:
        pil_img = page.render(scale=scale).to_pil()
        buf = io.BytesIO()
        if target == "png":
            pil_img.save(buf, format="PNG", optimize=True)
        else:
            pil_img.convert("RGB").save(buf, format="JPEG", quality=88)
        out.append(buf.getvalue())
    return out


# --- Default starter template ----------------------------------------------
#
# Generated programmatically rather than committed as a binary so it stays
# in sync with the documented token list — when we add a new supported
# token to the resolver registry, we add a line here too. Squadrons can
# download this, open it in PowerPoint, drop their logo + restyle, and
# re-upload as their custom template.

def generate_default_template() -> bytes:
    """Return a multi-slide .pptx populated with every supported token.

    Slide layout:
      1. Cover (mission name, theater, date/time)
      2. Flight & comms (per-flight callsign, aircraft, freq, TACAN)
      3. Weather
      4. Notes / free text
    """
    from pptx import Presentation as Prs
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Prs()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)
    LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
    ACCENT = RGBColor(0xFF, 0xA5, 0x00)
    DIM = RGBColor(0xAA, 0xAA, 0xAA)

    def add_bg(slide):
        # Fill background by adding a full-slide rectangle behind everything.
        from pptx.enum.shapes import MSO_SHAPE
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = DARK_BG
        rect.line.fill.background()
        # Move to back so other shapes layer on top.
        spTree = rect._element.getparent()
        spTree.remove(rect._element)
        spTree.insert(2, rect._element)

    def add_text(slide, x, y, w, h, text, size=18, bold=False, color=LIGHT, align_center=False):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align_center:
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Arial"
        return tf

    # ---------------- Slide 1: Cover ---------------------------------------
    s1 = prs.slides.add_slide(blank)
    add_bg(s1)
    add_text(s1, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5),
             "MISSION BRIEF", size=14, bold=True, color=ACCENT)
    add_text(s1, Inches(0.6), Inches(1.0), Inches(12), Inches(1.4),
             "{{mission.sortie}}", size=44, bold=True, color=LIGHT)
    add_text(s1, Inches(0.6), Inches(2.6), Inches(12), Inches(0.6),
             "{{mission.theater}}  ·  {{mission.date}}  ·  {{mission.time_zulu}}",
             size=20, color=DIM)
    add_text(s1, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
             "Lead: {{flight[0].callsign}}  ({{flight[0].aircraft}})",
             size=14, color=DIM)

    # ---------------- Slide 2: Flight & comms ------------------------------
    s2 = prs.slides.add_slide(blank)
    add_bg(s2)
    add_text(s2, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "FLIGHTS & COMMS", size=24, bold=True, color=ACCENT)

    # Two-flight table — squadrons typically brief 2–4 flights
    rows = [
        ["Callsign", "Aircraft", "Freq (MHz)", "TACAN", "ICLS"],
        ["{{flight[0].callsign}}", "{{flight[0].aircraft}}",
         "{{flight[0].frequency}}", "{{flight[0].tacan}}", "{{flight[0].icls}}"],
        ["{{flight[1].callsign}}", "{{flight[1].aircraft}}",
         "{{flight[1].frequency}}", "{{flight[1].tacan}}", "{{flight[1].icls}}"],
        ["{{flight[2].callsign}}", "{{flight[2].aircraft}}",
         "{{flight[2].frequency}}", "{{flight[2].tacan}}", "{{flight[2].icls}}"],
    ]
    table_shape = s2.shapes.add_table(
        rows=len(rows), cols=len(rows[0]),
        left=Inches(0.6), top=Inches(1.4), width=Inches(12), height=Inches(2.5),
    )
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.name = "Arial"
                    run.font.color.rgb = ACCENT if r_idx == 0 else LIGHT
                    if r_idx == 0:
                        run.font.bold = True

    # ---------------- Slide 3: Weather -------------------------------------
    s3 = prs.slides.add_slide(blank)
    add_bg(s3)
    add_text(s3, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "WEATHER", size=24, bold=True, color=ACCENT)
    wx_lines = [
        ("Cloud preset",   "{{weather.cloud_preset}}"),
        ("Visibility (m)", "{{weather.visibility_m}}"),
        ("QNH (inHg)",     "{{weather.qnh_inhg}}"),
        ("QNH (hPa)",      "{{weather.qnh_hpa}}"),
        ("Temperature",    "{{weather.temp_c}} °C"),
        ("Wind — surface", "{{weather.wind_surface}}"),
        ("Wind — 2000 ft", "{{weather.wind_2000}}"),
        ("Wind — 8000 ft", "{{weather.wind_8000}}"),
    ]
    for i, (label, value) in enumerate(wx_lines):
        y = Inches(1.4 + i * 0.55)
        add_text(s3, Inches(0.6), y, Inches(3), Inches(0.5), label, size=14, color=DIM)
        add_text(s3, Inches(3.8), y, Inches(8), Inches(0.5), value, size=16, color=LIGHT)

    # ---------------- Slide 4: Notes / free text ---------------------------
    s4 = prs.slides.add_slide(blank)
    add_bg(s4)
    add_text(s4, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "MISSION DESCRIPTION", size=24, bold=True, color=ACCENT)
    add_text(s4, Inches(0.6), Inches(1.4), Inches(12), Inches(5),
             "{{mission.description}}", size=14, color=LIGHT)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# --- Wing brief renderer (no template, direct .pptx generation) ----------
#
# Builds a presentation-ready .pptx from a WingBrief dict. Each section
# becomes one slide. Layout is consistent — dark-grey background,
# orange section headers, Arial body — matching the squadron-brief
# aesthetic of the kneeboard cards. Mission makers can re-style after
# download.

def _hex_is_dark(hexval: str) -> bool:
    """Perceptual-luminance check for a 'RRGGBB' hex string. True = dark."""
    try:
        h = hexval.strip().lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        # Rec.601 luma, 0-255.
        return (0.299 * r + 0.587 * g + 0.114 * b) < 128
    except Exception:
        return True


def _master_bg_is_dark(prs) -> bool:
    """Decide whether the uploaded template's slide-master background is
    dark. Drives the brief's text palette so content stays readable.

    Reads the master's <p:bg> fill: a literal srgbClr is luminance-tested;
    a schemeClr name is mapped through the master's clrMap and judged by
    its theme-slot prefix (lt* = light, dk*/tx* = dark). When the master
    has no explicit background we assume light (PowerPoint's default is a
    white slide), since that's the common 'branded blank' case and a wrong
    'dark' guess there would render invisible light text on white.
    """
    try:
        from pptx.oxml.ns import qn
        melem = prs.slide_masters[0].element
        bg = melem.find(".//" + qn("p:bg"))
        if bg is None:
            return False  # no explicit bg → PPT default white → light
        srgb = bg.find(".//" + qn("a:srgbClr"))
        if srgb is not None and srgb.get("val"):
            return _hex_is_dark(srgb.get("val"))
        scheme = bg.find(".//" + qn("a:schemeClr"))
        if scheme is None or not scheme.get("val"):
            return False
        name = scheme.get("val")
        clrmap = melem.find(".//" + qn("p:clrMap"))
        mapping = dict(clrmap.attrib) if clrmap is not None else {}
        slot = mapping.get(name, name)  # e.g. bg1 -> lt1
        if slot.startswith("lt"):
            return False
        if slot.startswith("dk") or slot.startswith("tx"):
            return True
        return False
    except Exception:
        # Inconclusive — assume light so we don't hide text on a white deck.
        return False


def _set_tf_text(tf, value: str) -> None:
    """Replace a text-frame's text while keeping the first run's formatting
    (font/size/color). Falls back to tf.text when there are no runs."""
    paras = tf.paragraphs
    if paras and paras[0].runs:
        paras[0].runs[0].text = value
        for r in paras[0].runs[1:]:
            r.text = ""
        for p in paras[1:]:
            for r in p.runs:
                r.text = ""
    else:
        tf.text = value


def _fill_template_cover(slide, title: str, subtitle: str) -> None:
    """Fill a template cover slide's title/subtitle. Handles both literal
    'Title'/'Sub Title' placeholder TEXT (Google-Slides-style exports) and
    real PowerPoint title/subtitle placeholders."""
    title_done = sub_done = False
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        norm = " ".join(sh.text_frame.text.split()).lower()
        if norm in ("title", "{{title}}", "[title]", "mission title") and title:
            _set_tf_text(sh.text_frame, title); title_done = True
        elif norm in ("sub title", "subtitle", "{{subtitle}}", "[subtitle]") and subtitle:
            _set_tf_text(sh.text_frame, subtitle); sub_done = True
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        for ph in slide.placeholders:
            t = ph.placeholder_format.type
            if not title_done and t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and title:
                _set_tf_text(ph.text_frame, title); title_done = True
            elif not sub_done and t == PP_PLACEHOLDER.SUBTITLE and subtitle:
                _set_tf_text(ph.text_frame, subtitle); sub_done = True
    except Exception:
        pass


def _draw_popup_mini_profile(slide, p: Dict[str, Any], *, x, y, w, h,
                              accent, dim, light, border) -> None:
    """Draw a thumbnail side-profile of one popup-attack profile inside
    (x, y, w, h) of the slide. Pure shape-based — no images — so it
    survives PowerPoint round-trips cleanly.

    Geometry mirrors utils/popupAttack.ts so the slide thumbnail matches
    the kneeboard card's chart. Reference points: IP / AP / PDP / RP /
    TGT / REC, drawn as small ovals with a single-letter label centred
    in each.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt, Emu

    FT_PER_NM = 6076.115
    INGRESS_DISPLAY_NM = 5.0
    atk = str(p.get("attackType") or "type1")

    def _f(k, default=0.0):
        v = p.get(k)
        try:
            return float(v) if v is not None else float(default)
        except (TypeError, ValueError):
            return float(default)

    t_elev   = _f("targetElevationFt", 100)
    vip_nm   = max(0.0, _f("vipDistanceNm", 8))
    pop_msl  = _f("popupAltitudeFtMsl", 8000)
    pop_ang  = max(1.0, _f("popupAngleDeg", 40))
    dive_ang = max(1.0, _f("diveAngleDeg", 30))
    rel_agl  = _f("releaseAltitudeFtAgl", 2000)
    ing_agl  = _f("ingressAltitudeFtAgl", 500)
    rec_agl  = _f("recoveryAltitudeFtAgl", ing_agl)
    ingress_msl = t_elev + ing_agl
    release_msl = t_elev + rel_agl
    recovery_msl = t_elev + rec_agl

    # Build (along-track NM, alt MSL) waypoints per attack type. Mirrors
    # the JS computePopupAttack in utils/popupAttack.ts — kept short:
    # only the points the chart needs (IP, AP, PDP/RP/TGT/REC as relevant).
    pts: List[tuple] = []  # (label, dist_nm, alt_msl)
    pts.append(("IP", 0.0, ingress_msl))
    ap_dist = INGRESS_DISPLAY_NM
    pts.append(("AP", ap_dist, ingress_msl))

    if atk == "laydown":
        tgt_dist = ap_dist + vip_nm
        pts.append(("RP",  tgt_dist, release_msl))
        pts.append(("TGT", tgt_dist, t_elev))
        pts.append(("REC", tgt_dist + 1.5, recovery_msl))
    elif atk == "loft":
        climb_v = pop_msl - ingress_msl
        climb_h_nm = max(0.0, (climb_v / max(1e-3, math.tan(math.radians(pop_ang)))) / FT_PER_NM)
        rp_dist = ap_dist + climb_h_nm
        pts.append(("RP", rp_dist, pop_msl))
        tgt_dist = rp_dist + max(0.5, vip_nm * 0.4)
        pts.append(("TGT", tgt_dist, t_elev))
        pts.append(("REC", tgt_dist + 1.5, recovery_msl))
    elif atk == "dive":
        dive_v = ingress_msl - release_msl
        dive_h_nm = max(0.0, (dive_v / max(1e-3, math.tan(math.radians(dive_ang)))) / FT_PER_NM)
        rp_dist = ap_dist + dive_h_nm
        pts.append(("RP", rp_dist, release_msl))
        tgt_dist = rp_dist + 0.5
        pts.append(("TGT", tgt_dist, t_elev))
        pts.append(("REC", tgt_dist + 1.5, recovery_msl))
    else:
        # type1/2/3 popup — same math
        climb_v = pop_msl - ingress_msl
        climb_h_nm = max(0.0, (climb_v / max(1e-3, math.tan(math.radians(pop_ang)))) / FT_PER_NM)
        pdp_dist = ap_dist + climb_h_nm
        pts.append(("PDP", pdp_dist, pop_msl))
        dive_v = pop_msl - release_msl
        dive_h_nm = max(0.0, (dive_v / max(1e-3, math.tan(math.radians(dive_ang)))) / FT_PER_NM)
        rp_dist = pdp_dist + dive_h_nm
        pts.append(("RP", rp_dist, release_msl))
        tgt_dist = rp_dist + 0.5
        pts.append(("TGT", tgt_dist, t_elev))
        pts.append(("REC", tgt_dist + 1.5, recovery_msl))

    # Scale into the chart box. Leave a 2 px padding margin at the edges.
    x0, y0, ww, hh = int(x), int(y), int(w), int(h)
    pad = Pt(3)
    xmin, xmax = 0.0, max(pt[1] for pt in pts) or 1.0
    ymin = min(min(pt[2] for pt in pts), t_elev)
    ymax = max(pt[2] for pt in pts)
    ymax = max(ymax, ymin + 1.0)  # avoid divide-by-zero on degenerate profiles
    inner_w = ww - 2 * int(pad)
    inner_h = hh - 2 * int(pad)

    def xy(d_nm, alt_msl):
        u = (d_nm - xmin) / (xmax - xmin) if xmax > xmin else 0
        v = 1 - (alt_msl - ymin) / (ymax - ymin) if ymax > ymin else 1
        return x0 + int(pad) + int(u * inner_w), y0 + int(pad) + int(v * inner_h)

    # Chart border (subtle, for visual containment).
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y0, ww, hh)
    box.fill.background()  # transparent fill
    box.line.color.rgb = border; box.line.width = Pt(0.5)

    # Ground line (target elevation).
    gx1, gy = xy(xmin, t_elev)
    gx2, _ = xy(xmax, t_elev)
    ground = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx1, gy, gx2 - gx1, Emu(6000))
    ground.fill.solid(); ground.fill.fore_color.rgb = dim
    ground.line.fill.background()

    # Trajectory segments — straight lines between non-TGT points.
    traj = [pt for pt in pts if pt[0] != "TGT"]
    line_thick = Emu(8000)
    for a, b in zip(traj, traj[1:]):
        ax, ay = xy(a[1], a[2])
        bx, by = xy(b[1], b[2])
        # Use a rotated thin rectangle as a line segment. Math: position +
        # rotation between (ax, ay) and (bx, by). python-pptx's rotation
        # is in degrees clockwise; centre the rect on the segment midpoint.
        dx, dy = bx - ax, by - ay
        length = max(1, int(math.hypot(dx, dy)))
        angle = math.degrees(math.atan2(dy, dx))
        mx, my = (ax + bx) // 2, (ay + by) // 2
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, mx - length // 2, my - line_thick // 2,
            length, line_thick,
        )
        seg.fill.solid(); seg.fill.fore_color.rgb = accent
        seg.line.fill.background()
        seg.rotation = angle

    # Markers — small ovals with a 1-letter label.
    dot_r = Pt(4)
    for label, d_nm, alt_msl in pts:
        cx, cy = xy(d_nm, alt_msl)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx - int(dot_r), cy - int(dot_r),
            int(dot_r * 2), int(dot_r * 2),
        )
        dot.fill.solid(); dot.fill.fore_color.rgb = accent
        dot.line.color.rgb = light; dot.line.width = Pt(0.5)
        # Label letter — first character (IP/AP/PDP/RP/TGT/REC → I/A/P/R/T/R).
        lbl = "T" if label == "TGT" else label[0]
        tx = slide.shapes.add_textbox(
            cx + int(dot_r) + Pt(2), cy - Pt(6),
            Pt(20), Pt(12),
        )
        tf = tx.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        run = para.add_run(); run.text = lbl
        run.font.size = Pt(7); run.font.bold = True; run.font.color.rgb = light


def _slide_has_content(slide) -> bool:
    """True if a slide carries its own text or picture (master-inherited
    branding doesn't count — it's on every slide)."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return True
        if getattr(sh, "shape_type", None) == 13:  # PICTURE
            return True
    return False


def render_wing_brief(brief: Dict[str, Any], base_template_b64: Optional[str] = None,
                      top_margin_in: Optional[float] = None) -> bytes:
    """Render a WingBrief dict to .pptx bytes.

    base_template_b64: optional base64 of a squadron .pptx. When given,
    the brief is built ON that template — its own slide(s) + theme are
    preserved as the cover/branding and the content slides are appended;
    the built-in cover is dropped. Any failure falls back to the default
    deck so the brief always renders.

    Slide order (Phase 1 — wing brief only):
      1. Cover
      2. Theatre overview
      3. Scenario
      4. Commander's intent
      5. Threats (table)
      6. Force composition (flights table)
      7. Comms (kv list)
      8. Mission flow
      9. Timeline (table)
     10. Notes / special instructions
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import base64 as _b64

    # Optional base template (v0.9.79). Build ON the uploaded .pptx so its
    # slides + theme survive; fall back to the default deck on any error.
    use_template = False
    if base_template_b64:
        try:
            _tb = _b64.b64decode(base_template_b64.split(",", 1)[-1])
            prs = Presentation(io.BytesIO(_tb))
            use_template = True
        except Exception:
            prs = Presentation()
    else:
        prs = Presentation()
    if not use_template:
        prs.slide_width = Inches(13.333)   # 16:9
        prs.slide_height = Inches(7.5)

    # How many slides the template already had — its branding cover(s).
    n_template_slides = len(prs.slides._sldIdLst)

    # Blank layout — custom templates may not have the standard 7-layout
    # set, so prefer a placeholder-free layout and fall back gracefully.
    def _pick_blank_layout():
        for _lay in prs.slide_layouts:
            try:
                if len(_lay.placeholders) == 0:
                    return _lay
            except Exception:
                continue
        try:
            return prs.slide_layouts[6]
        except Exception:
            return prs.slide_layouts[0]
    BLANK = _pick_blank_layout()

    # Palette (v0.9.80). Default = dark (built-in deck). When building on a
    # user template we DON'T paint our own background — the template's
    # master branding shows through — so the text palette must match the
    # template's background brightness, or text vanishes (e.g. light grey
    # on a white master). Auto-detect dark vs light from the master bg.
    dark = True if not use_template else _master_bg_is_dark(prs)
    if dark:
        BG = RGBColor(0x1A, 0x1A, 0x1A)
        LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
        BRIGHT = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT = RGBColor(0xFF, 0xA5, 0x00)
        DIM = RGBColor(0xAA, 0xAA, 0xAA)
        BORDER = RGBColor(0x55, 0x55, 0x55)
        TABLE_HEADER_BG = RGBColor(0x33, 0x33, 0x33)
        CELL_BG = RGBColor(0x1A, 0x1A, 0x1A)
    else:
        # Light template — dark text, light table panels, a strong dark
        # amber accent that reads on white.
        BG = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT = RGBColor(0x1A, 0x1A, 0x1A)
        BRIGHT = RGBColor(0x00, 0x00, 0x00)
        ACCENT = RGBColor(0xB8, 0x74, 0x0C)
        DIM = RGBColor(0x55, 0x55, 0x55)
        BORDER = RGBColor(0x99, 0x99, 0x99)
        TABLE_HEADER_BG = RGBColor(0xD8, 0xD8, 0xD8)
        CELL_BG = RGBColor(0xF3, 0xF3, 0xF3)

    # Content top inset (v0.9.81) — only when building on a template, so
    # section headers + content drop below the template's branded header
    # band / logos instead of colliding with them. Caller-tunable (clamped
    # 0-4"); default 1.2".
    if use_template:
        _top = 1.2 if top_margin_in is None else max(0.0, min(float(top_margin_in), 4.0))
    else:
        _top = 0.0
    _MY = Inches(_top)

    # Paginated-table budgets (v0.9.83). When a template top-margin pushes
    # content down, shrink the per-page row counts + table height so tables
    # keep a ~0.5" bottom margin instead of running off the slide. Default
    # deck (no margin) keeps its existing, tuned counts.
    if _top > 0:
        _tbl_h_in = max(2.6, 7.5 - (1.4 + _top) - 0.5)
        _flights_per_slide = max(4, int((_tbl_h_in - 0.4) / 0.46))
        _threats_per_slide = max(3, int((_tbl_h_in - 0.7) / 0.65))
    else:
        _tbl_h_in = 5.5
        _flights_per_slide = 12
        _threats_per_slide = 7

    # ---------- helpers ---------------------------------------------------

    def _apply_bg(slide):
        # On a user template, let the master's branding/background show.
        if use_template:
            return
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = BG
        rect.line.fill.background()
        spTree = rect._element.getparent()
        spTree.remove(rect._element)
        spTree.insert(2, rect._element)

    def _txt(slide, x, y, w, h, text, *, size=18, bold=False, color=LIGHT,
             align_center=False, italic=False):
        # Shift content down by the template inset; trim height so a tall
        # body box doesn't run off the bottom edge.
        tx = slide.shapes.add_textbox(x, y + _MY, w, max(h - _MY, Inches(0.4)))
        tf = tx.text_frame
        tf.word_wrap = True
        # Multi-paragraph support — split body text on \n into separate paragraphs.
        lines = text.split("\n") if text else [""]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align_center:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.name = "Arial"
        return tf

    def _slide_header(slide, label):
        """Top-of-slide section title in accent color."""
        _txt(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             label, size=24, bold=True, color=ACCENT)
        # Thin underline
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05) + _MY,
                                      Inches(12.1), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

    def _table(slide, x, y, w, h, headers, rows, col_widths=None):
        ncols = len(headers)
        nrows = len(rows) + 1  # +1 for header row
        shape = slide.shapes.add_table(nrows, ncols, x, y + _MY, w, h)
        table = shape.table
        # Column widths
        if col_widths:
            for i, cw in enumerate(col_widths):
                table.columns[i].width = cw

        # Header row
        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HEADER_BG
            cell.text = header
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.color.rgb = ACCENT
                    r.font.size = Pt(13)
                    r.font.name = "Arial"

        # Body rows
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = CELL_BG
                cell.text = str(val) if val is not None else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = LIGHT
                        r.font.size = Pt(12)
                        r.font.name = "Arial"
        return table

    # ---------- Slide 1: Cover -------------------------------------------
    # Layout (16:9, 13.333" × 7.5"):
    #   * Optional cover hero image — fills upper ~45% of slide as
    #     full-width banner. Title and metadata float below.
    #   * Without cover image: top half is a solid accent bar with a
    #     "WING BRIEF" eyebrow; the layout still feels intentional.
    #   * Mission name centred, dramatic 56pt bold, in the lower band.
    #   * Theater · Date · Time line below the title in the accent colour.
    #   * Optional squadron logo overlays top-right regardless of layout.
    #   * Bottom accent bar with a subtle "BRIEF · v1" eyebrow on the
    #     right and a left/right pair of vertical accent ticks.
    import base64

    s = prs.slides.add_slide(BLANK); _apply_bg(s)

    # Helper: decode a base64 image (with optional data: prefix) into a
    # BytesIO stream, or None if invalid. Used for both the cover hero
    # and the squadron logo.
    def _decode_image(b64: str):
        b = (b64 or "").strip()
        if not b:
            return None
        try:
            if "," in b and b.lstrip().startswith("data:"):
                b = b.split(",", 1)[1]
            return io.BytesIO(base64.b64decode(b, validate=False))
        except Exception:
            return None

    cover_img = _decode_image(brief.get("cover_image_base64") or "")
    HERO_HEIGHT = Inches(3.6)  # upper 48% of the slide

    if cover_img:
        # Hero image — full slide width, top of slide. python-pptx will
        # crop/letterbox via the size we specify; we set both width and
        # height so the image fills the banner edge-to-edge.
        try:
            s.shapes.add_picture(
                cover_img, 0, 0,
                width=prs.slide_width, height=HERO_HEIGHT,
            )
            # Subtle dark overlay at the bottom of the hero so the
            # transition into the dark slide bg isn't a hard line.
            grad = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, HERO_HEIGHT - Inches(0.6),
                prs.slide_width, Inches(0.6),
            )
            grad.fill.solid(); grad.fill.fore_color.rgb = BG
            # Soft fade by setting transparency through the XML — easier
            # to just set a flat 50% opaque dark bar than to add a real
            # gradient stop in this codepath.
            grad.fill.fore_color.rgb = BG
            grad.line.fill.background()
        except Exception:
            cover_img = None  # decode-but-render failure: fall back to text-only

    if not cover_img:
        # No hero image — give the top a styled accent block so the
        # cover doesn't feel empty. Two horizontal bars frame the
        # eyebrow text.
        bar_top = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08),
        )
        bar_top.fill.solid(); bar_top.fill.fore_color.rgb = ACCENT
        bar_top.line.fill.background()

        # Vertical accent ticks at left + right edges — quiet visual frame
        for x_in in (Inches(0.4), prs.slide_width - Inches(0.5)):
            tick = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_in, Inches(1.4), Inches(0.06), Inches(2.5),
            )
            tick.fill.solid(); tick.fill.fore_color.rgb = ACCENT
            tick.line.fill.background()

    # Squadron logo — rendered top-right, overlays cover image when present
    logo_img = _decode_image(brief.get("logo_base64") or "")
    if logo_img:
        try:
            s.shapes.add_picture(
                logo_img, Inches(11.5), Inches(0.3),
                height=Inches(1.4),
            )
        except Exception:
            pass

    # Eyebrow ("WING BRIEF") — top-left, smaller when there's a hero so
    # it doesn't fight the image, larger otherwise
    eyebrow_y = Inches(0.4) if cover_img else Inches(0.55)
    _txt(s, Inches(0.6), eyebrow_y, Inches(8), Inches(0.5),
         "WING BRIEF", size=14, bold=True, color=ACCENT)

    # Title block — sits below the hero or in the styled top half
    title_top = HERO_HEIGHT + Inches(0.4) if cover_img else Inches(2.6)

    _txt(s, Inches(0.6), title_top, Inches(12.1), Inches(1.4),
         brief["mission_name"], size=56, bold=True, color=BRIGHT,
         align_center=True)

    # Theater · Date · Time strip — accent colour, smaller
    sub_top = title_top + Inches(1.5)
    _txt(s, Inches(0.6), sub_top, Inches(12.1), Inches(0.5),
         f"{brief['theater'].upper()}   ·   {brief['date']}   ·   "
         f"TAKEOFF {brief['time_zulu']}",
         size=18, bold=True, color=ACCENT, align_center=True)

    # Bottom accent bar — visual anchor at the bottom of the slide
    bottom_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.08),
        prs.slide_width, Inches(0.08),
    )
    bottom_bar.fill.solid(); bottom_bar.fill.fore_color.rgb = ACCENT
    bottom_bar.line.fill.background()

    # ---------- Slide 2: Theatre overview -------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "THEATRE OVERVIEW")
    _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
         brief["theatre_overview"], size=16, color=LIGHT)

    # ---------- Slide 3: Scenario ----------------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "SCENARIO")
    _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
         brief["scenario"], size=16, color=LIGHT)

    # ---------- Slide 4: Commander's intent ------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "COMMANDER'S INTENT")
    _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
         brief["commanders_intent"], size=16, color=LIGHT)

    # ---------- Slide 4b: Route overview map (client-rendered) -----------
    # All flight tracks + threat rings on one image (captureOverviewImage,
    # 760x520). Fitted by that aspect so it never distorts. Optional.
    overview_img = _decode_image(brief.get("route_overview_base64") or "")
    if overview_img:
        s = prs.slides.add_slide(BLANK); _apply_bg(s)
        _slide_header(s, "ROUTE OVERVIEW")
        iw, ih = 760, 520
        avail_w, avail_h = Inches(12.1), Inches(5.7)
        scale = min(avail_w / iw, avail_h / ih)
        w = int(iw * scale); h = int(ih * scale)
        x = int((prs.slide_width - w) / 2); y = Inches(1.45)
        try:
            s.shapes.add_picture(overview_img, x, y, width=w, height=h)
        except Exception:
            pass

    # ---------- Slide 5: Threats ----------------------------------------
    # Spatially-clustered threat areas, ranked by tier then engagement
    # range. Pilots see "STRATEGIC / TACTICAL / SHORAD / AAA" at a glance
    # before drilling into individual systems. Slide overflows to a
    # second slide automatically when the cluster count exceeds what
    # fits readably (~10 rows at this font size).
    threats_list = brief.get("threats") or []

    # Tier color hints on the slide — keep palette restrained
    TIER_COLOR = {
        "STRATEGIC": RGBColor(0xFF, 0x55, 0x55),  # red — long-range area defence
        "TACTICAL":  RGBColor(0xFF, 0xA5, 0x00),  # accent orange — medium SAM
        "MIXED":     RGBColor(0xFF, 0xA5, 0x00),  # treat as tactical-leaning
        "SHORAD":    RGBColor(0xFF, 0xD9, 0x66),  # amber — short-range
        "MANPAD":    RGBColor(0xCC, 0xCC, 0x88),  # tan — IR threats
        "AAA":       RGBColor(0xAA, 0xAA, 0xAA),  # neutral grey — gun-only
        "OTHER":     RGBColor(0x88, 0x88, 0x88),  # dim — unrecognised
    }

    def _render_threats_slide(rows_for_slide, page_idx, total_pages):
        s = prs.slides.add_slide(BLANK); _apply_bg(s)
        title = "SURFACE THREATS" if total_pages == 1 else f"SURFACE THREATS ({page_idx}/{total_pages})"
        _slide_header(s, title)

        if page_idx == 1:
            # Lead-line on first page only
            tier_count = sum(1 for r in threats_list)
            _txt(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.4),
                 f"{tier_count} threat area(s) — position is bearing/range (nm) "
                 f"from bullseye; WEZ is max engagement range in cluster.",
                 size=11, color=DIM, italic=True)

        # Layout — column ranges are non-overlapping and total 12.1" wide:
        #   TIER         x 0.6  w 1.5  → ends 2.1
        #   COMPOSITION  x 2.2  w 5.6  → ends 7.8
        #   POSITION     x 7.9  w 2.0  → ends 9.9  (BE bearing/range)
        #   WEZ          x 10.0 w 2.7  → ends 12.7
        # Was overlapping because COMPOSITION (6.0" wide) often had long
        # text like "1× SA-11 + 4× SA-15 Tor M1 + 6× ZU-23 Emplacement"
        # that wrapped to a second line and bled into the next row's space
        # since ROW_H stayed 0.55". Fixed by reducing rows-per-slide,
        # shrinking the composition font, raising row height.
        TOP = Inches(1.7) if page_idx == 1 else Inches(1.3)
        ROW_H = Inches(0.65)  # bumped from 0.55 to allow 2-line wrap headroom
        COL_TIER = (Inches(0.6),  Inches(1.5))
        COL_COMP = (Inches(2.2),  Inches(5.6))
        COL_POS  = (Inches(7.9),  Inches(2.0))
        COL_WEZ  = (Inches(10.0), Inches(2.7))

        # Header band
        for label, (x, w) in (("TIER", COL_TIER), ("COMPOSITION", COL_COMP),
                              ("POSITION", COL_POS), ("WEZ", COL_WEZ)):
            _txt(s, x, TOP, w, ROW_H, label, size=12, bold=True, color=ACCENT)
        # Underline below the column-header row. Must include the template
        # top-margin (_MY) like the _txt/_table helpers — without it this
        # rule floated up to the top of the slide and read as a stray grey
        # line under a template's header band. Uses the palette BORDER so
        # it adapts to light templates instead of a hardcoded dark grey.
        underline = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), TOP + Inches(0.5) + _MY,
            Inches(12.1), Inches(0.025),
        )
        underline.fill.solid(); underline.fill.fore_color.rgb = BORDER
        underline.line.fill.background()

        for i, t in enumerate(rows_for_slide):
            y = TOP + Inches(0.7) + ROW_H * i
            tier = t.get("tier") or "OTHER"
            tier_color = TIER_COLOR.get(tier, LIGHT)
            wez_km = t.get("range_km", 0)
            wez_nm = t.get("range_nm", 0)

            # Composition can be long; shrink font when content suggests
            # it'd wrap. Threshold roughly matches 5.6" at 13pt Arial.
            comp = t.get("composition") or t.get("name") or "?"
            comp_size = 13 if len(comp) <= 60 else 11

            _txt(s, COL_TIER[0], y, COL_TIER[1], ROW_H,
                 tier, size=13, bold=True, color=tier_color)
            _txt(s, COL_COMP[0], y, COL_COMP[1], ROW_H,
                 comp, size=comp_size, color=LIGHT)
            _txt(s, COL_POS[0], y, COL_POS[1], ROW_H,
                 t.get("location", "—"),
                 size=13, color=LIGHT, bold=True)
            _txt(s, COL_WEZ[0], y, COL_WEZ[1], ROW_H,
                 f"{wez_nm:.0f} nm  ({wez_km:.0f} km)" if wez_nm > 0 else "—",
                 size=12, color=LIGHT)

    # ---------- Slide 4c: Threat brief (AI prose, optional) --------------
    # Short paragraph that frames the threat picture for the lead before the
    # detailed surface/air tables. Empty by default — only renders when the
    # planner / AI fills it in. (v1.15.x)
    tn = (brief.get("threat_narrative") or "").strip()
    if tn:
        s = prs.slides.add_slide(BLANK); _apply_bg(s)
        _slide_header(s, "THREAT BRIEF")
        _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
             tn, size=18, color=LIGHT)

    if threats_list:
        # 7 rows × 0.65" = 4.55" of rows + 0.7" header headroom = 5.25"
        # plus the 1.7" header zone = 6.95" total — fits in 7.5" slide.
        # Was 9 rows × 0.55" = 4.95" but with text wrapping that exceeded
        # the row band and overlapped the next row's text.
        ROWS_PER_SLIDE = _threats_per_slide
        n = len(threats_list)
        total_pages = (n + ROWS_PER_SLIDE - 1) // ROWS_PER_SLIDE
        for page_idx in range(1, total_pages + 1):
            start = (page_idx - 1) * ROWS_PER_SLIDE
            _render_threats_slide(threats_list[start:start + ROWS_PER_SLIDE],
                                  page_idx, total_pages)
    else:
        s = prs.slides.add_slide(BLANK); _apply_bg(s)
        _slide_header(s, "SURFACE THREATS")
        msg = "No surface (SAM/AAA) threats detected in this mission."
        if brief.get("air_threats"):
            msg += "  Enemy air picture follows on the AIR THREATS slide."
        _txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1),
             msg, size=18, color=DIM, italic=True)

    # ---------- Slide 5b: Air threats -----------------------------------
    # Enemy aircraft aggregated by airframe TYPE (e.g. "8× Su-27"), each with a
    # capability rundown for friendly pilots: class, A2A weapons/WEZ, and how to
    # fight it. Paginated like the surface threats slide.
    air_list = brief.get("air_threats") or []
    if air_list:
        AIR_COLS = (  # (label, x, w) — non-overlapping, total 12.1"
            ("TYPE",    Inches(0.6),  Inches(2.3)),
            ("CLASS",   Inches(3.0),  Inches(2.4)),
            ("WEAPONS", Inches(5.5),  Inches(3.3)),
            ("NOTES",   Inches(8.9),  Inches(3.8)),
        )
        A_ROWS = _threats_per_slide
        n_air = len(air_list)
        air_pages = (n_air + A_ROWS - 1) // A_ROWS
        for pidx in range(air_pages):
            s = prs.slides.add_slide(BLANK); _apply_bg(s)
            a_title = "AIR THREATS" if air_pages == 1 else f"AIR THREATS ({pidx + 1}/{air_pages})"
            _slide_header(s, a_title)
            if pidx == 0:
                _txt(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.4),
                     f"{n_air} enemy airframe type(s) on the map — A2A weapons and "
                     f"how to fight each.",
                     size=11, color=DIM, italic=True)
            A_TOP = Inches(1.7) if pidx == 0 else Inches(1.3)
            A_ROW_H = Inches(0.7)  # taller — NOTES wraps to ~2 lines
            for label, x, w in AIR_COLS:
                _txt(s, x, A_TOP, w, A_ROW_H, label, size=12, bold=True, color=ACCENT)
            a_underline = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), A_TOP + Inches(0.5) + _MY,
                Inches(12.1), Inches(0.025),
            )
            a_underline.fill.solid(); a_underline.fill.fore_color.rgb = BORDER
            a_underline.line.fill.background()
            for i, a in enumerate(air_list[pidx * A_ROWS:(pidx + 1) * A_ROWS]):
                y = A_TOP + Inches(0.7) + A_ROW_H * i
                _txt(s, AIR_COLS[0][1], y, AIR_COLS[0][2], A_ROW_H,
                     a.get("composition") or "?", size=14, bold=True, color=ACCENT)
                _txt(s, AIR_COLS[1][1], y, AIR_COLS[1][2], A_ROW_H,
                     a.get("airframe_class") or "—", size=12, color=LIGHT)
                _txt(s, AIR_COLS[2][1], y, AIR_COLS[2][2], A_ROW_H,
                     a.get("weapons") or "—", size=11, color=LIGHT)
                _txt(s, AIR_COLS[3][1], y, AIR_COLS[3][2], A_ROW_H,
                     a.get("notes") or "—", size=10, color=DIM)

    # ---------- Slide 6: Force composition -------------------------------
    # Paginated when there are more flights than fit comfortably on one
    # slide. Each row is ~0.4" tall at 11pt Arial; 5.5" of rows = ~13
    # rows safely. Was overflowing on missions with 14+ flights because
    # python-pptx tables grow downward past their requested height when
    # row count exceeds what fits — causing rows to fall off the slide.
    flights_list = brief.get("flights") or []
    if flights_list:
        FLIGHTS_PER_SLIDE = _flights_per_slide
        n = len(flights_list)
        flight_pages = (n + FLIGHTS_PER_SLIDE - 1) // FLIGHTS_PER_SLIDE
        for page_idx in range(1, flight_pages + 1):
            start = (page_idx - 1) * FLIGHTS_PER_SLIDE
            chunk = flights_list[start:start + FLIGHTS_PER_SLIDE]
            s = prs.slides.add_slide(BLANK); _apply_bg(s)
            title = ("FRIENDLY FORCES" if flight_pages == 1
                     else f"FRIENDLY FORCES ({page_idx}/{flight_pages})")
            _slide_header(s, title)
            _table(
                s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(_tbl_h_in),
                ["Callsign", "Aircraft", "#", "Role", "Freq (MHz)", "TACAN", "Home Plate"],
                [[f.get("callsign", ""), f.get("aircraft", ""), str(f.get("count", "")),
                  f.get("role", ""), f.get("frequency", ""), f.get("tacan", ""),
                  f.get("home_plate", "")] for f in chunk],
                col_widths=[Inches(2.0), Inches(2.5), Inches(0.5), Inches(1.5),
                            Inches(1.5), Inches(1.2), Inches(2.5)],
            )
    else:
        s = prs.slides.add_slide(BLANK); _apply_bg(s)
        _slide_header(s, "FRIENDLY FORCES")
        _txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1),
             "No player flights detected.", size=18, color=DIM, italic=True)

    # ---------- Slide 7: Comms ------------------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "COMMS")
    for i, c in enumerate(brief["comms"]):
        y = Inches(1.4 + i * 0.5)
        _txt(s, Inches(0.6), y, Inches(3.5), Inches(0.5), c.get("label", ""),
             size=15, color=DIM)
        _txt(s, Inches(4.5), y, Inches(8), Inches(0.5), c.get("value", ""),
             size=17, color=LIGHT, bold=True)

    # ---------- Slide 8: Mission flow -----------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "MISSION FLOW")
    _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
         brief["mission_flow"], size=15, color=LIGHT)

    # ---------- Slide 9: Timeline ---------------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "TIMELINE")
    if brief["timeline"]:
        _table(
            s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(max(2.5, 5.0 - _top)),
            ["Phase", "Time (Z)", "Note"],
            [[r.get("phase", ""), r.get("time_zulu", ""), r.get("note", "")]
             for r in brief["timeline"]],
            col_widths=[Inches(3.0), Inches(2.0), Inches(7.1)],
        )

    # ---------- Slide 10: Notes -----------------------------------------
    s = prs.slides.add_slide(BLANK); _apply_bg(s)
    _slide_header(s, "SPECIAL INSTRUCTIONS / NOTES")
    notes_text = brief.get("notes") or (
        "Edit this section to add ROE, special procedures, contingency "
        "plans, fragments of code-words, divert decisions, etc."
    )
    is_placeholder = not brief.get("notes")
    _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
         notes_text, size=15, color=DIM if is_placeholder else LIGHT,
         italic=is_placeholder)

    # ---------- Slide 11: Popup attack profiles (optional) ---------------
    # Pages through the kneeboard popup-attack profiles, one row per
    # profile per slide (max 4 per slide so each gets enough vertical
    # room for the parameter ladder). Skipped entirely when the planner
    # didn't define any. (v1.17.6)
    popups = brief.get("popup_attacks") or []
    if popups:
        ATTACK_LABELS = {
            "type1":   "Type 1 Popup",
            "type2":   "Type 2 Popup",
            "type3":   "Type 3 Popup",
            "laydown": "Lay-Down",
            "loft":    "Loft (Toss)",
            "dive":    "Straight Dive",
        }
        ROWS_PER_SLIDE = 4
        n_pop = len(popups)
        pop_pages = (n_pop + ROWS_PER_SLIDE - 1) // ROWS_PER_SLIDE
        for pidx in range(pop_pages):
            s = prs.slides.add_slide(BLANK); _apply_bg(s)
            head = "POPUP ATTACK" if pop_pages == 1 else f"POPUP ATTACK ({pidx + 1}/{pop_pages})"
            _slide_header(s, head)
            if pidx == 0:
                _txt(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.4),
                     f"{n_pop} attack profile(s) defined for this mission — "
                     "Apply per-flight as required.",
                     size=11, color=DIM, italic=True)
            P_TOP = Inches(1.7) if pidx == 0 else Inches(1.3)
            P_ROW_H = Inches(1.35)
            for j, p in enumerate(popups[pidx * ROWS_PER_SLIDE:(pidx + 1) * ROWS_PER_SLIDE]):
                y = P_TOP + P_ROW_H * j
                atk = str(p.get("attackType") or "type1")
                atk_label = ATTACK_LABELS.get(atk, atk)
                name = str(p.get("name") or f"Attack {pidx * ROWS_PER_SLIDE + j + 1}")
                # Profile name + type chip
                _txt(s, Inches(0.6), y, Inches(5.0), Inches(0.4),
                     name, size=15, bold=True, color=BRIGHT)
                _txt(s, Inches(5.6), y + Inches(0.05), Inches(2.4), Inches(0.4),
                     atk_label, size=12, bold=True, color=ACCENT)
                # Parameters — two columns of three rows each.
                def _p(k, default=""):
                    v = p.get(k)
                    return v if v is not None else default
                t_elev   = f"{int(_p('targetElevationFt') or 0):,} ft MSL"
                vip      = f"{_p('vipDistanceNm')} NM"
                pop_alt  = f"{int(_p('popupAltitudeFtMsl') or 0):,} ft MSL"
                pop_ang  = f"{_p('popupAngleDeg')}°"
                div_ang  = f"{_p('diveAngleDeg')}°"
                off      = f"{_p('angleOffsetDeg')}°"
                rel_alt  = f"{int(_p('releaseAltitudeFtAgl') or 0):,} ft AGL"
                rel_spd  = f"{_p('releaseSpeedKts')} kt"
                ing_alt  = f"{int(_p('ingressAltitudeFtAgl') or 0):,} ft AGL"
                ing_spd  = f"{_p('ingressSpeedKts')} kt"
                left_col = (
                    f"TGT elev:  {t_elev}\n"
                    f"VIP dist:  {vip}\n"
                    f"Popup:     {pop_alt}  /  {pop_ang}\n"
                    f"Offset:    {off}"
                )
                right_col = (
                    f"Dive ang:  {div_ang}\n"
                    f"Release:   {rel_alt}  /  {rel_spd}\n"
                    f"Ingress:   {ing_alt}  /  {ing_spd}\n"
                )
                _txt(s, Inches(0.6),  y + Inches(0.42), Inches(4.4), Inches(0.95),
                     left_col, size=12, color=LIGHT)
                _txt(s, Inches(5.1),  y + Inches(0.42), Inches(4.4), Inches(0.95),
                     right_col, size=12, color=LIGHT)
                # Mini side-profile chart — pilots want to SEE the geometry,
                # not just read parameters. Drawn with plain pptx shapes so
                # it round-trips through PowerPoint cleanly (no embedded
                # SVG / images).
                _draw_popup_mini_profile(
                    s, p, x=Inches(9.6), y=y + Inches(0.4),
                    w=Inches(3.5), h=Inches(0.9),
                    accent=ACCENT, dim=DIM, light=LIGHT, border=BORDER,
                )
                # Row separator
                if j < min(ROWS_PER_SLIDE, n_pop - pidx * ROWS_PER_SLIDE) - 1:
                    sep = s.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0.6), y + P_ROW_H - Inches(0.05),
                        Inches(12.1), Inches(0.012),
                    )
                    sep.fill.solid(); sep.fill.fore_color.rgb = BORDER
                    sep.line.fill.background()

    # Template post-processing (v0.9.84):
    #   1. Drop the built-in cover (the slide WE added right after the
    #      template's own slides) — the template's slide(s) are the cover.
    #   2. Fill the template's first slide ('Title' / 'Sub Title') with the
    #      mission name + theatre/date/time.
    #   3. Drop any blank template slides after the cover (e.g. an empty
    #      second slide) so they don't show as blank pages.
    if use_template:
        # 1) remove built-in cover
        try:
            sld_id_lst = prs.slides._sldIdLst
            ids = list(sld_id_lst)
            if len(ids) > n_template_slides:
                sld_id_lst.remove(ids[n_template_slides])
        except Exception:
            pass
        # 2) fill the cover (first template slide)
        try:
            cover_title = str(brief.get("mission_name") or "")
            sub_parts = [
                str(brief.get("theater") or "").upper(),
                str(brief.get("date") or ""),
                (f"TAKEOFF {brief.get('time_zulu')}" if brief.get("time_zulu") else ""),
            ]
            cover_sub = "   ·   ".join(p for p in sub_parts if p)
            if len(prs.slides._sldIdLst) > 0:
                _fill_template_cover(prs.slides[0], cover_title, cover_sub)
        except Exception:
            pass
        # 3) drop blank template slides after the cover (keep slide 0)
        try:
            ids = list(prs.slides._sldIdLst)
            limit = min(n_template_slides, len(ids))
            for i in range(limit - 1, 0, -1):
                if not _slide_has_content(prs.slides[i]):
                    prs.slides._sldIdLst.remove(ids[i])
        except Exception:
            pass

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def render_flight_brief(brief: Dict[str, Any]) -> bytes:
    """Render one FlightBrief dict to a compact 4-5 slide .pptx.

    Slide order (per-flight, ~4 slides total):
      1. Cover (callsign / aircraft / mission / time)
      2. Tasking (free-text, mission-type-aware default)
      3. Route (waypoint table)
      4. Comms + Fuel (per-flight freq/TACAN/ICLS + joker/bingo/RTB)
      5. Notes (only if non-empty — keeps short briefs tight)

    Sized 16:9 to match the wing brief, same dark/orange palette so
    a wing brief and per-flight brief look like one product when
    bundled together.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    BG = RGBColor(0x1A, 0x1A, 0x1A)
    LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
    BRIGHT = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = RGBColor(0xFF, 0xA5, 0x00)
    DIM = RGBColor(0xAA, 0xAA, 0xAA)
    TABLE_HEADER_BG = RGBColor(0x33, 0x33, 0x33)

    def _bg(slide):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                       prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = BG
        rect.line.fill.background()
        spTree = rect._element.getparent()
        spTree.remove(rect._element); spTree.insert(2, rect._element)

    def _txt(slide, x, y, w, h, text, *, size=18, bold=False, color=LIGHT,
             align_center=False, italic=False):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame; tf.word_wrap = True
        lines = text.split("\n") if text else [""]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align_center:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color; r.font.name = "Arial"
        return tf

    def _header(slide, label):
        _txt(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             label, size=24, bold=True, color=ACCENT)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.6), Inches(1.05),
                                       Inches(12.1), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

    def _table(slide, x, y, w, h, headers, rows, col_widths=None):
        ncols = len(headers); nrows = len(rows) + 1
        shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
        table = shape.table
        if col_widths:
            for i, cw in enumerate(col_widths):
                table.columns[i].width = cw
        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HEADER_BG
            cell.text = header
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True; r.font.color.rgb = ACCENT
                    r.font.size = Pt(13); r.font.name = "Arial"
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = BG
                cell.text = str(val) if val is not None else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = LIGHT
                        r.font.size = Pt(11); r.font.name = "Arial"

    # ---------- 1: Cover ----------------------------------------------
    s = prs.slides.add_slide(BLANK); _bg(s)
    _txt(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
         f"FLIGHT BRIEF — {brief.get('callsign', '')}", size=14, bold=True,
         color=ACCENT, align_center=True)
    _txt(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
         brief.get("callsign", ""), size=56, bold=True, color=BRIGHT,
         align_center=True)
    _txt(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.5),
         f"{brief.get('count', '?')}× {brief.get('aircraft', '?')}  ·  "
         f"Role: {brief.get('role', '')}",
         size=18, color=DIM, align_center=True)
    _txt(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
         f"Mission: {brief.get('mission_name', '')}",
         size=14, color=DIM, align_center=True)
    _txt(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
         f"{brief.get('theater', '')}  ·  {brief.get('date', '')}  ·  "
         f"Takeoff {brief.get('time_zulu', '')}",
         size=14, color=DIM, align_center=True)

    # ---------- 2: Tasking --------------------------------------------
    s = prs.slides.add_slide(BLANK); _bg(s)
    _header(s, "TASKING")
    _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
         brief.get("tasking", ""), size=18, color=LIGHT)

    # ---------- 3: Route ----------------------------------------------
    s = prs.slides.add_slide(BLANK); _bg(s)
    _header(s, "ROUTE")
    wps = brief.get("waypoints") or []
    if wps:
        _table(
            s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4),
            ["#", "Name", "Alt (ft MSL)", "Speed (kt)", "Distance (nm)", "ETA (Z)"],
            [[w.get("number", ""), w.get("name", ""),
              w.get("altitude_ft", ""), w.get("speed_kt", ""),
              w.get("distance_nm", ""), w.get("eta_zulu", "")]
             for w in wps],
            col_widths=[Inches(0.7), Inches(3.5), Inches(2.0),
                        Inches(2.0), Inches(2.0), Inches(1.9)],
        )
    else:
        _txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1),
             "No waypoints in mission file.", size=18, color=DIM, italic=True)

    # ---------- 3b: Route map (client-rendered PNG, optional) ---------
    rm = (brief.get("route_map_base64") or "").strip()
    if rm:
        try:
            import base64
            if rm.startswith("data:") and "," in rm:
                rm = rm.split(",", 1)[1]
            img_bytes = base64.b64decode(rm, validate=False)
            iw, ih = 560, 400  # captureRouteImage default aspect (7:5)
            try:
                from PIL import Image as _PImg
                with _PImg.open(io.BytesIO(img_bytes)) as _im:
                    iw, ih = _im.size
            except Exception:
                pass
            s = prs.slides.add_slide(BLANK); _bg(s)
            _header(s, "ROUTE MAP")
            avail_w, avail_h = Inches(12.1), Inches(5.6)
            scale = min(avail_w / iw, avail_h / ih)
            w = int(iw * scale); h = int(ih * scale)
            x = int((prs.slide_width - w) / 2); y = Inches(1.45)
            s.shapes.add_picture(io.BytesIO(img_bytes), x, y, width=w, height=h)
        except Exception:
            pass

    # ---------- 4: Comms + Fuel ---------------------------------------
    s = prs.slides.add_slide(BLANK); _bg(s)
    _header(s, "COMMS + FUEL")
    # Two columns side by side
    col1_x = Inches(0.6); col1_w = Inches(6.0)
    col2_x = Inches(7.0); col2_w = Inches(5.7)

    _txt(s, col1_x, Inches(1.4), col1_w, Inches(0.4),
         "COMMS", size=14, bold=True, color=ACCENT)
    comms_rows = [
        ("Primary",   f"{brief.get('frequency', '')} MHz" if brief.get('frequency') else ""),
        ("TACAN",     brief.get("tacan", "")),
        ("ICLS",      brief.get("icls", "")),
        ("Home Plate", brief.get("home_plate", "")),
        ("Divert",    brief.get("divert", "")),
    ]
    for i, (label, value) in enumerate(comms_rows):
        y = Inches(1.9 + i * 0.5)
        _txt(s, col1_x, y, Inches(2.0), Inches(0.4), label, size=13, color=DIM)
        _txt(s, col1_x + Inches(2.2), y, Inches(3.5), Inches(0.4),
             value or "—", size=15, color=LIGHT, bold=True)

    _txt(s, col2_x, Inches(1.4), col2_w, Inches(0.4),
         "FUEL LADDER", size=14, bold=True, color=ACCENT)
    fuel_rows = [
        ("Joker",  f"{brief.get('fuel_joker_lbs', 0):,} lbs"),
        ("Bingo",  f"{brief.get('fuel_bingo_lbs', 0):,} lbs"),
        ("RTB",    f"{brief.get('fuel_rtb_lbs', 0):,} lbs"),
    ]
    for i, (label, value) in enumerate(fuel_rows):
        y = Inches(1.9 + i * 0.55)
        _txt(s, col2_x, y, Inches(2.0), Inches(0.4), label, size=13, color=DIM)
        _txt(s, col2_x + Inches(2.2), y, Inches(3.5), Inches(0.4),
             value, size=18, color=LIGHT, bold=True)
    _txt(s, col2_x, Inches(4.0), col2_w, Inches(0.4),
         "(Edit fuel values in the editor — defaults are placeholders)",
         size=10, color=DIM, italic=True)

    # This flight's own schedule (package timeline lives on the wing brief).
    tl = brief.get("timeline") or []
    if tl:
        _txt(s, col1_x, Inches(4.7), Inches(12.1), Inches(0.4),
             "SCHEDULE", size=14, bold=True, color=ACCENT)
        for i, row in enumerate(tl[:5]):
            y = Inches(5.2 + i * 0.42)
            _txt(s, col1_x, y, Inches(2.0), Inches(0.4),
                 row.get("phase", ""), size=13, color=DIM)
            _txt(s, col1_x + Inches(2.1), y, Inches(1.6), Inches(0.4),
                 row.get("time_zulu", ""), size=15, bold=True, color=LIGHT)
            _txt(s, col1_x + Inches(3.9), y, Inches(7.8), Inches(0.4),
                 row.get("note", ""), size=12, color=DIM)

    # ---------- 5: Notes (only if non-empty) ---------------------------
    notes = brief.get("notes") or ""
    if notes.strip():
        s = prs.slides.add_slide(BLANK); _bg(s)
        _header(s, "NOTES / SPECIAL INSTRUCTIONS")
        _txt(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8),
             notes, size=15, color=LIGHT)

    # ---------- 6: Popup attack appendix (only if profiles exist) -------
    # Compact per-flight version of the wing-brief slide — same row layout,
    # same mini side-profile sketches, paged 4 profiles/slide. v1.17.8.
    flight_pops = brief.get("popup_attacks") or []
    if flight_pops:
        ATTACK_LABELS = {
            "type1":   "Type 1 Popup",
            "type2":   "Type 2 Popup",
            "type3":   "Type 3 Popup",
            "laydown": "Lay-Down",
            "loft":    "Loft (Toss)",
            "dive":    "Straight Dive",
        }
        ROWS_PER_SLIDE = 4
        n_pop = len(flight_pops)
        pop_pages = (n_pop + ROWS_PER_SLIDE - 1) // ROWS_PER_SLIDE
        BORDER = RGBColor(0x55, 0x55, 0x55)
        for pidx in range(pop_pages):
            s = prs.slides.add_slide(BLANK); _bg(s)
            head_label = "POPUP ATTACK" if pop_pages == 1 else f"POPUP ATTACK ({pidx + 1}/{pop_pages})"
            _header(s, head_label)
            P_TOP = Inches(1.4)
            P_ROW_H = Inches(1.35)
            page_slice = flight_pops[pidx * ROWS_PER_SLIDE:(pidx + 1) * ROWS_PER_SLIDE]
            for j, p in enumerate(page_slice):
                y = P_TOP + P_ROW_H * j
                atk = str(p.get("attackType") or "type1")
                atk_label = ATTACK_LABELS.get(atk, atk)
                name = str(p.get("name") or f"Attack {pidx * ROWS_PER_SLIDE + j + 1}")
                _txt(s, Inches(0.6), y, Inches(5.0), Inches(0.4),
                     name, size=15, bold=True, color=BRIGHT)
                _txt(s, Inches(5.6), y + Inches(0.05), Inches(2.4), Inches(0.4),
                     atk_label, size=12, bold=True, color=ACCENT)
                def _pf(k, default=""):
                    v = p.get(k)
                    return v if v is not None else default
                left_col = (
                    f"TGT elev:  {int(_pf('targetElevationFt') or 0):,} ft MSL\n"
                    f"VIP dist:  {_pf('vipDistanceNm')} NM\n"
                    f"Popup:     {int(_pf('popupAltitudeFtMsl') or 0):,} ft MSL  /  {_pf('popupAngleDeg')}°\n"
                    f"Offset:    {_pf('angleOffsetDeg')}°"
                )
                right_col = (
                    f"Dive ang:  {_pf('diveAngleDeg')}°\n"
                    f"Release:   {int(_pf('releaseAltitudeFtAgl') or 0):,} ft AGL  /  {_pf('releaseSpeedKts')} kt\n"
                    f"Ingress:   {int(_pf('ingressAltitudeFtAgl') or 0):,} ft AGL  /  {_pf('ingressSpeedKts')} kt\n"
                )
                _txt(s, Inches(0.6),  y + Inches(0.42), Inches(4.4), Inches(0.95),
                     left_col, size=12, color=LIGHT)
                _txt(s, Inches(5.1),  y + Inches(0.42), Inches(4.4), Inches(0.95),
                     right_col, size=12, color=LIGHT)
                _draw_popup_mini_profile(
                    s, p, x=Inches(9.6), y=y + Inches(0.4),
                    w=Inches(3.5), h=Inches(0.9),
                    accent=ACCENT, dim=DIM, light=LIGHT, border=BORDER,
                )
                if j < len(page_slice) - 1:
                    sep = s.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0.6), y + P_ROW_H - Inches(0.05),
                        Inches(12.1), Inches(0.012),
                    )
                    sep.fill.solid(); sep.fill.fore_color.rgb = BORDER
                    sep.line.fill.background()

    out = io.BytesIO(); prs.save(out)
    return out.getvalue()


def _rasterize_pdf_via_ghostscript(
    pdf_bytes: bytes, target: Literal["png", "jpg"], dpi: int,
) -> List[bytes]:
    """Ghostscript fallback for environments without pypdfium2."""
    gs = shutil.which("gs") or shutil.which("gswin64c") or shutil.which("gswin32c")
    if not gs:
        raise RuntimeError(
            "Neither pypdfium2 nor Ghostscript found. "
            "Install pypdfium2 (`pip install pypdfium2`) for image export."
        )
    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        pdf_path = tmpdir / "in.pdf"
        pdf_path.write_bytes(pdf_bytes)
        out_pattern = tmpdir / f"page_%03d.{target}"
        device = "png16m" if target == "png" else "jpeg"
        cmd = [
            gs, "-q", "-dNOPAUSE", "-dBATCH", "-dSAFER",
            f"-sDEVICE={device}", f"-r{dpi}",
            f"-sOutputFile={out_pattern}",
            str(pdf_path),
        ]
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
        images = sorted(tmpdir.glob(f"page_*.{target}"))
        return [p.read_bytes() for p in images]
