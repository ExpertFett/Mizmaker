"""Render a PKT intel packet (services.pkt_builder output) to .pptx bytes.

A classified intel-document look: crimson TOP SECRET banners top and bottom of
every card, a TN identifier + declassification date, and recognition cards that
pair the in-house silhouette with how-to-fight bullets. Same python-pptx
approach as the brief renderer; art (dark silhouettes) lives in
backend/assets/aircraft/dark/.
"""
from __future__ import annotations

import io
import os
from typing import Any, Dict, List, Optional

_DARK_SIL_DIR = os.path.join(os.path.dirname(__file__), "..", "assets", "aircraft", "dark")
_dark_cache: Dict[str, Optional[bytes]] = {}


def _dark_sil(family: Optional[str]) -> Optional[bytes]:
    if not family:
        return None
    if family in _dark_cache:
        return _dark_cache[family]
    data = None
    try:
        p = os.path.join(_DARK_SIL_DIR, f"{family}.png")
        if os.path.isfile(p):
            with open(p, "rb") as f:
                data = f.read()
    except Exception:
        data = None
    _dark_cache[family] = data
    return data


def render_pkt(pkt: Dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # Intel-document palette
    PAPER = RGBColor(0xF3, 0xF1, 0xEA)
    INK = RGBColor(0x20, 0x24, 0x2B)
    CRIMSON = RGBColor(0x8A, 0x12, 0x20)
    DIM = RGBColor(0x6A, 0x70, 0x78)
    BORDER = RGBColor(0xCF, 0xCC, 0xC0)
    CELL = RGBColor(0xFB, 0xFA, 0xF6)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DISPLAY, BODY, MONO, LABEL = "Oswald", "Barlow", "Consolas", "Barlow Semi Condensed"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    marking = pkt.get("marking") or "TOP SECRET // REL TO COALITION"
    decl = pkt.get("decl_on") or ""
    _tn = [4050]

    def next_tn():
        _tn[0] += 1
        return f"TN {_tn[0]}"

    def txt(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
            font=BODY, align=None, spacing=None):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate((text or "").split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align is not None:
                p.alignment = align
            if spacing is not None:
                p.line_spacing = spacing
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
        return tf

    def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if fill is None:
            s.fill.background()
        else:
            s.fill.solid(); s.fill.fore_color.rgb = fill
        if line is None:
            s.line.fill.background()
        else:
            s.line.color.rgb = line; s.line.width = Pt(line_w)
        return s

    def new_card(section, *, tn=True):
        s = prs.slides.add_slide(BLANK)
        rect(s, 0, 0, W, H, fill=PAPER)  # paper
        # Classified banners
        for y in (0, H - Inches(0.32)):
            rect(s, 0, y, W, Inches(0.32), fill=CRIMSON)
            txt(s, 0, y + Inches(0.03), W, Inches(0.28), marking, size=13,
                bold=True, color=WHITE, font=LABEL, align=PP_ALIGN.CENTER)
        # Section label + identifier + rule
        txt(s, Inches(0.5), Inches(0.46), Inches(8), Inches(0.5), section,
            size=20, bold=True, color=CRIMSON, font=DISPLAY)
        if tn:
            ident = next_tn()
            line2 = f"{ident}\n{('DECL ON: ' + decl) if decl else ''}".strip()
            txt(s, Inches(9.0), Inches(0.44), Inches(3.8), Inches(0.6), line2,
                size=12, color=DIM, font=MONO, align=PP_ALIGN.RIGHT)
        rect(s, Inches(0.5), Inches(0.92), Inches(12.33), Inches(0.03), fill=INK)
        return s

    def add_sil(slide, family, x, y, h):
        data = _dark_sil(family)
        if not data:
            return False
        try:
            slide.shapes.add_picture(io.BytesIO(data), x, y, height=h)
            return True
        except Exception:
            return False

    # ---------------- Cover ----------------
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=INK)
    for y in (0, H - Inches(0.36)):
        rect(s, 0, y, W, Inches(0.36), fill=CRIMSON)
        txt(s, 0, y + Inches(0.04), W, Inches(0.3), marking, size=14, bold=True,
            color=WHITE, font=LABEL, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
        "INTELLIGENCE PACKET", size=22, bold=True, color=RGBColor(0xC9, 0xA2, 0x4A),
        font=LABEL)
    txt(s, Inches(0.66), Inches(2.8), Inches(12), Inches(1.6),
        pkt.get("mission_name") or "MISSION", size=54, bold=True, color=WHITE,
        font=DISPLAY)
    meta = f"{pkt.get('theater','')}   ·   {pkt.get('date','')}   ·   {pkt.get('packet_id','')}"
    txt(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5), meta, size=18,
        color=RGBColor(0xC7, 0xCC, 0xD1), font=BODY)
    if decl:
        txt(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
            f"DECLASSIFY ON: {decl}  ·  FICTION / TRAINING USE ONLY", size=12,
            color=DIM, font=MONO)

    # ---------------- Friendly Forces ----------------
    fr = pkt.get("friendly") or {}
    flights = fr.get("flights") or []
    s = new_card("FRIENDLY FORCES")
    txt(s, Inches(0.5), Inches(1.05), Inches(7.4), Inches(0.5),
        f"Order of Battle · {('CV ' + fr['carrier']) if fr.get('carrier') else 'Package'}",
        size=26, bold=True, color=INK, font=DISPLAY)
    # OOB table (left)
    headers = ["CALLSIGN", "AIRCRAFT", "ROLE", "FREQ"]
    rows = [[f.get("callsign", ""), f"{f.get('aircraft','')} ×{f.get('count','')}",
             f.get("role", ""), f.get("frequency", "")] for f in flights[:12]]
    if rows:
        n = len(rows) + 1
        tbl = s.shapes.add_table(n, 4, Inches(0.5), Inches(1.75),
                                 Inches(6.3), Inches(0.34 * n)).table
        for w, ci in zip((Inches(1.9), Inches(2.0), Inches(1.4), Inches(1.0)), range(4)):
            tbl.columns[ci].width = w
        for ci, hdr in enumerate(headers):
            c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = INK
            c.text = hdr
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True; r.font.size = Pt(11); r.font.name = LABEL
                    r.font.color.rgb = PAPER
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci); c.fill.solid(); c.fill.fore_color.rgb = CELL
                c.text = str(val)
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(12); r.font.name = BODY; r.font.color.rgb = INK
    else:
        txt(s, Inches(0.5), Inches(1.8), Inches(6.3), Inches(0.5),
            "No player flights detected in mission.", size=15, color=DIM, font=BODY)
    # Theatre map (right)
    ao = pkt.get("ao_center")
    box_x, box_y, box_w, box_h = Inches(7.1), Inches(1.55), Inches(5.73), Inches(5.05)
    rect(s, box_x, box_y, box_w, box_h, fill=None, line=INK, line_w=2.5)
    if isinstance(ao, dict):
        try:
            from services import ao_imagery
            png = ao_imagery.fetch_ao_image(float(ao["lat"]), float(ao["lon"]),
                                            float(ao.get("span_km") or 300),
                                            1280, 1130, brightness=0.95, saturate=0.7)
            if png:
                s.shapes.add_picture(io.BytesIO(png), box_x + Inches(0.04),
                                     box_y + Inches(0.04), width=box_w - Inches(0.08),
                                     height=box_h - Inches(0.08))
        except Exception:
            pass
    txt(s, box_x, box_y + box_h - Inches(0.4), box_w, Inches(0.34),
        f"  AO · {pkt.get('theater','')}", size=12, color=WHITE, font=MONO)

    # ---------------- A/A Threat grid ----------------
    air = pkt.get("air_threats") or []
    if air:
        s = new_card("CURRENT A/A THREAT")
        txt(s, Inches(0.5), Inches(1.02), Inches(11), Inches(0.5),
            "Enemy Air Order of Battle", size=24, bold=True, color=INK, font=DISPLAY)
        cols = 4
        cw, ch, gap = Inches(3.02), Inches(2.0), Inches(0.12)
        x0, y0 = Inches(0.5), Inches(1.65)
        for i, a in enumerate(air[:8]):
            cx = x0 + (cw + gap) * (i % cols)
            cy = y0 + (ch + gap) * (i // cols)
            rect(s, cx, cy, cw, ch, fill=CELL, line=BORDER)
            add_sil(s, a.get("silhouette"), cx + cw / 2 - Inches(0.55),
                    cy + Inches(0.12), Inches(1.05))
            txt(s, cx, cy + Inches(1.28), cw, Inches(0.4),
                f"{a.get('name','')}   ×{a.get('count','')}", size=17, bold=True,
                color=INK, font=DISPLAY, align=PP_ALIGN.CENTER)
            txt(s, cx, cy + Inches(1.66), cw, Inches(0.3),
                f"{a.get('reporting','')}", size=12, color=DIM, font=LABEL,
                align=PP_ALIGN.CENTER)

    # ---------------- A/A Recognition (one per type) ----------------
    for a in air:
        s = new_card("A/A THREAT · RECOGNITION")
        box = Inches(0.5), Inches(1.35), Inches(4.6), Inches(4.9)
        rect(s, *box, fill=None, line=INK, line_w=2)
        add_sil(s, a.get("silhouette"), Inches(1.1), Inches(1.7), Inches(4.2))
        tx = Inches(5.5)
        txt(s, tx, Inches(1.2), Inches(7.3), Inches(1.0), a.get("name", ""),
            size=48, bold=True, color=INK, font=DISPLAY)
        rep = a.get("reporting", ""); role = a.get("role", "")
        txt(s, tx, Inches(2.15), Inches(7.3), Inches(0.5),
            (f'"{rep}"' if rep else "") + (f" · {role}" if role else ""),
            size=22, bold=True, color=CRIMSON, font=DISPLAY)
        bullets = "\n".join("•  " + b for b in (a.get("ttps") or []))
        txt(s, tx, Inches(2.95), Inches(7.4), Inches(3.4), bullets, size=17,
            color=INK, font=BODY, spacing=1.25)

    # ---------------- Surface-Threat Recognition ----------------
    for sfc in (pkt.get("surface_threats") or []):
        s = new_card("SURFACE THREAT · RECOGNITION")
        txt(s, Inches(0.5), Inches(1.2), Inches(9), Inches(1.0), sfc.get("name", ""),
            size=48, bold=True, color=INK, font=DISPLAY)
        tier = sfc.get("tier", ""); wez = sfc.get("wez")
        sub = tier + (f"  ·  WEZ {wez:.0f} NM" if isinstance(wez, (int, float)) and wez else "")
        txt(s, Inches(0.52), Inches(2.15), Inches(9), Inches(0.5), sub, size=22,
            bold=True, color=CRIMSON, font=DISPLAY)
        if sfc.get("composition"):
            txt(s, Inches(0.52), Inches(2.7), Inches(9), Inches(0.4),
                sfc["composition"], size=14, color=DIM, font=BODY)
        bullets = "\n".join("•  " + b for b in (sfc.get("tactics") or []))
        txt(s, Inches(0.52), Inches(3.35), Inches(11.8), Inches(3.0), bullets,
            size=18, color=INK, font=BODY, spacing=1.3)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
