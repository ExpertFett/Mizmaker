"""Speaker-notes injection in services/brief_renderer.

Tests that:
1. A brief without speaker_notes renders exactly as before — no notes
   are added to any slide.
2. A brief WITH speaker_notes attaches each note to the slide whose
   section-header text matches the corresponding key.
3. The cover slide (which has no section header) gets the "cover" note
   unconditionally when one is provided.
4. Unknown keys in the speaker_notes map are ignored silently (no crash,
   no garbage on slides).
"""

from __future__ import annotations

import io

import pytest


@pytest.fixture
def minimal_brief() -> dict:
    """A brief with every required field so render_wing_brief reaches the
    end-of-function speaker-notes pass."""
    return {
        "mission_name": "Test Mission",
        "theater": "Caucasus",
        "date": "2026-06-06",
        "time_zulu": "0900Z",
        "coalition": "blue",
        "theatre_overview": "Test theatre overview.",
        "scenario": "Test scenario.",
        "commanders_intent": "Test intent.",
        "mission_flow": "Test flow.",
        "notes": "Test notes.",
        "timeline": [],
        "threats": [],
        "air_threats": [],
        "flights": [],
        "comms": [],
        "logo_base64": "",
        "cover_image_base64": "",
        "threat_narrative": "",
        "popup_attacks": [],
    }


def _read_notes(pptx_bytes: bytes) -> list[tuple[str, str]]:
    """Return [(first-header-text, speaker-notes-text), ...] for every
    slide in the rendered deck. Header text is the first non-empty run
    on the slide; speaker notes is the full notes_text_frame content."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    out: list[tuple[str, str]] = []
    for slide in prs.slides:
        header = ""
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    t = (run.text or "").strip()
                    if t:
                        header = t
                        break
                if header:
                    break
            if header:
                break
        note = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text or ""
        out.append((header, note.strip()))
    return out


def test_brief_without_speaker_notes_has_no_notes(minimal_brief):
    """Baseline: when speaker_notes is absent, no slide should carry any
    speaker-notes text."""
    pytest.importorskip("pptx")
    from services.brief_renderer import render_wing_brief
    out_bytes = render_wing_brief(minimal_brief)
    for _, note in _read_notes(out_bytes):
        assert note == "", f"expected no notes, got: {note!r}"


def test_speaker_notes_attach_to_matching_section_headers(minimal_brief):
    """When speaker_notes provides intent + threats + notes, those keys
    should land on the slides whose headers say so."""
    pytest.importorskip("pptx")
    from services.brief_renderer import render_wing_brief
    minimal_brief["speaker_notes"] = {
        "cover": "Cover talking point.",
        "intent": "Intent talking point.",
        "notes": "Notes talking point.",
    }
    out_bytes = render_wing_brief(minimal_brief)
    rows = _read_notes(out_bytes)
    # Cover is the first slide and has no section-header check — it gets
    # the "cover" note unconditionally.
    assert rows[0][1] == "Cover talking point.", rows[0]
    # Find the intent + notes slides by header substring (header is in
    # SECTION HEADER caps; match case-insensitively).
    by_header: dict[str, str] = {h.lower(): n for h, n in rows}
    intent_match = next(
        (n for h, n in by_header.items() if "intent" in h),
        None,
    )
    assert intent_match == "Intent talking point."
    notes_match = next(
        (n for h, n in by_header.items() if "notes" in h or "instructions" in h),
        None,
    )
    assert notes_match == "Notes talking point."


def test_unknown_speaker_notes_keys_are_ignored(minimal_brief):
    """A bogus key in the speaker_notes map should be a silent no-op —
    no exception, no garbage on any slide."""
    pytest.importorskip("pptx")
    from services.brief_renderer import render_wing_brief
    minimal_brief["speaker_notes"] = {
        "unknown_key": "should never appear",
        "completely_made_up": "ditto",
    }
    out_bytes = render_wing_brief(minimal_brief)
    for _, note in _read_notes(out_bytes):
        assert "should never appear" not in note
        assert "ditto" not in note


def test_empty_string_note_is_a_noop(minimal_brief):
    """A note value of "" or "   " should be skipped — we shouldn't
    create an empty notes_slide for it."""
    pytest.importorskip("pptx")
    from services.brief_renderer import render_wing_brief
    minimal_brief["speaker_notes"] = {"intent": "   "}
    out_bytes = render_wing_brief(minimal_brief)
    for _, note in _read_notes(out_bytes):
        assert note == "" or note.strip() == ""
