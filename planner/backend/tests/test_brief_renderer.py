"""Tests for the PowerPoint brief renderer.

Generates a minimal .pptx in-memory per test rather than committing a binary
fixture — keeps the test self-contained and avoids drift between the fixture
file on disk and the assertions here.
"""
import io
import re

import pytest
from pptx import Presentation
from pptx.util import Inches

from services.brief_renderer import scan_template, render_template


def _make_template(*paragraphs: str) -> bytes:
    """Build a 1-slide .pptx where each arg is one body paragraph."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.text = paragraphs[0] if paragraphs else ""
    for p in paragraphs[1:]:
        tf.add_paragraph().text = p
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _slide_text(pptx_bytes: bytes) -> str:
    """Concatenate all paragraph text from slide 1 of a .pptx."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                parts.append("".join(r.text for r in p.runs))
    return "\n".join(parts)


# --- scan_template -----------------------------------------------------------

class TestScanTemplate:
    def test_finds_single_token(self):
        tokens = scan_template(_make_template("Theater: {{mission.theater}}"))
        assert tokens == ["mission.theater"]

    def test_finds_multiple_tokens_dedupes(self):
        tokens = scan_template(_make_template(
            "Mission: {{mission.name}} on {{mission.date}}",
            "Theater: {{mission.theater}}",
            "Repeat: {{mission.theater}}",
        ))
        assert tokens == sorted(["mission.name", "mission.date", "mission.theater"])

    def test_finds_array_indexed_tokens(self):
        tokens = scan_template(_make_template(
            "{{flight[0].callsign}} and {{flight[1].callsign}}"
        ))
        assert "flight[0].callsign" in tokens
        assert "flight[1].callsign" in tokens

    def test_no_tokens_in_plain_template(self):
        assert scan_template(_make_template("Just plain text, no tokens")) == []

    def test_rejects_malformed_pptx(self):
        with pytest.raises(ValueError):
            scan_template(b"not a pptx file")

    def test_ignores_invalid_token_syntax(self):
        # Only patterns matching {{name.with.dots}} are tokens. Stuff like
        # `{{ }}` or `{{1bad}}` (starts with digit) is left alone.
        tokens = scan_template(_make_template("{{ }} {{1bad}} {{good.token}}"))
        assert tokens == ["good.token"]


# --- render_template ---------------------------------------------------------

class TestRenderTemplate:
    def test_substitutes_known_token(self):
        tmpl = _make_template("Theater: {{mission.theater}}")
        rendered = render_template(tmpl, {"mission.theater": "Kola"})
        assert "Theater: Kola" in _slide_text(rendered)
        assert "{{mission.theater}}" not in _slide_text(rendered)

    def test_leaves_missing_token_as_literal(self):
        # The whole point of the soft-fail behaviour: user can spot what
        # didn't get filled in the rendered output.
        tmpl = _make_template("Set: {{mission.theater}}", "Unset: {{mission.callsign}}")
        rendered = render_template(tmpl, {"mission.theater": "Kola"})
        text = _slide_text(rendered)
        assert "Set: Kola" in text
        assert "{{mission.callsign}}" in text

    def test_substitutes_multiple_tokens_one_paragraph(self):
        tmpl = _make_template("{{a}} and {{b}} and {{c}}")
        rendered = render_template(tmpl, {"a": "X", "b": "Y", "c": "Z"})
        assert "X and Y and Z" in _slide_text(rendered)

    def test_handles_array_index_tokens(self):
        tmpl = _make_template("Lead: {{flight[0].callsign}}, Wing: {{flight[1].callsign}}")
        rendered = render_template(tmpl, {
            "flight[0].callsign": "BENGAL11",
            "flight[1].callsign": "BENGAL12",
        })
        assert "Lead: BENGAL11, Wing: BENGAL12" in _slide_text(rendered)

    def test_empty_values_dict_leaves_template_untouched(self):
        tmpl = _make_template("All: {{a}} {{b}} {{c}}")
        rendered = render_template(tmpl, {})
        assert "All: {{a}} {{b}} {{c}}" in _slide_text(rendered)

    def test_value_with_special_chars(self):
        # Substitution shouldn't mangle special chars like \n, %, &
        tmpl = _make_template("Notes: {{notes}}")
        rendered = render_template(tmpl, {"notes": "Line1 & Line2 — 100%"})
        assert "Line1 & Line2 — 100%" in _slide_text(rendered)

    def test_table_cell_substitution(self):
        # Tables are common in briefings (waypoint list, comm card).
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(
            rows=2, cols=2, left=Inches(1), top=Inches(1),
            width=Inches(6), height=Inches(2),
        )
        table = table_shape.table
        table.cell(0, 0).text = "Theater"
        table.cell(0, 1).text = "{{mission.theater}}"
        table.cell(1, 0).text = "Date"
        table.cell(1, 1).text = "{{mission.date}}"
        out = io.BytesIO(); prs.save(out)

        rendered = render_template(out.getvalue(), {
            "mission.theater": "Kola",
            "mission.date": "2026-04-25",
        })
        # Read back the table from the rendered file
        prs2 = Presentation(io.BytesIO(rendered))
        cells = [c.text for c in prs2.slides[0].shapes[0].table.iter_cells()]
        assert "Kola" in cells
        assert "2026-04-25" in cells

    def test_rejects_malformed_pptx(self):
        with pytest.raises(ValueError):
            render_template(b"definitely not a pptx", {"a": "b"})

    def test_render_round_trip_through_scan(self):
        """The output of render should still be a valid .pptx that scan
        can reopen — and any unfilled tokens should still be detectable."""
        tmpl = _make_template("Filled: {{a}}, Unfilled: {{b}}")
        rendered = render_template(tmpl, {"a": "DONE"})
        rescanned = scan_template(rendered)
        assert rescanned == ["b"]


# --- convert_pptx (PDF / PNG / JPG) ----------------------------------------
# Skipped automatically if LibreOffice isn't on PATH — these tests run on
# the CI image which installs libreoffice-impress, and on dev boxes that
# have it. They don't run on a stock Python install with no LibreOffice.

from services.brief_renderer import (
    convert_pptx, is_conversion_available, LibreOfficeNotFoundError,
)

skip_no_libreoffice = pytest.mark.skipif(
    not is_conversion_available(),
    reason="LibreOffice not installed — PDF/PNG/JPG conversion tests skipped",
)


class TestConvertPptx:
    def test_pptx_passthrough_no_libreoffice_required(self):
        """target='pptx' is a no-op pass-through and must work without soffice."""
        tmpl = _make_template("Hello")
        out, mime = convert_pptx(tmpl, "pptx")
        assert out == tmpl
        assert "presentation" in mime

    def test_unsupported_format_raises(self):
        with pytest.raises(ValueError):
            convert_pptx(_make_template("x"), "doc")  # type: ignore[arg-type]

    @skip_no_libreoffice
    def test_pdf_conversion(self):
        out, mime = convert_pptx(_make_template("PDF test"), "pdf")
        assert mime == "application/pdf"
        assert out.startswith(b"%PDF-")  # valid PDF magic

    @skip_no_libreoffice
    def test_png_conversion_returns_zip_of_slides(self):
        # 2-slide template so we can verify multiple images come through
        prs = Presentation()
        for i in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            tx.text_frame.text = f"Slide {i + 1}"
        buf = io.BytesIO(); prs.save(buf)

        out, mime = convert_pptx(buf.getvalue(), "png")
        assert mime == "application/zip"

        import zipfile
        with zipfile.ZipFile(io.BytesIO(out)) as zf:
            names = zf.namelist()
            assert len(names) == 2
            assert all(n.endswith(".png") for n in names)
            # Each entry should be a real PNG (magic bytes)
            for n in names:
                data = zf.read(n)
                assert data[:8] == b"\x89PNG\r\n\x1a\n"

    @skip_no_libreoffice
    def test_jpg_conversion_returns_zip_of_slides(self):
        out, mime = convert_pptx(_make_template("JPG test"), "jpg")
        assert mime == "application/zip"
        import zipfile
        with zipfile.ZipFile(io.BytesIO(out)) as zf:
            names = zf.namelist()
            assert len(names) >= 1
            assert all(n.endswith(".jpg") for n in names)
            data = zf.read(names[0])
            assert data[:3] == b"\xff\xd8\xff"  # JPEG magic

    def test_libreoffice_not_found_raises_specific_error(self, monkeypatch):
        """Force the not-found path even on a machine with soffice installed,
        so we can verify the error message points users to the install steps."""
        from services import brief_renderer
        monkeypatch.setattr(brief_renderer, "_find_soffice", lambda: None)
        with pytest.raises(LibreOfficeNotFoundError) as exc:
            convert_pptx(_make_template("x"), "pdf")
        assert "LibreOffice" in str(exc.value)
        assert "PowerPoint download (.pptx) still works" in str(exc.value)


# --- render_wing_brief: base-template support (v0.9.79) ----------------------

def _minimal_wing_brief() -> dict:
    """Smallest WingBrief dict render_wing_brief will accept."""
    return {
        "mission_name": "TEST OP", "theater": "Caucasus",
        "date": "2026-05-21", "time_zulu": "0830Z", "coalition": "blue",
        "logo_base64": "", "cover_image_base64": "",
        "theatre_overview": "Overview.", "scenario": "Scenario.",
        "commanders_intent": "Purpose: x\nMethod: y\nEnd State: z",
        "mission_flow": "1. Push", "notes": "",
        "timeline": [{"phase": "PUSH", "time_zulu": "0815Z", "note": "go"}],
        "threats": [{"name": "SA-11", "type": "SAM", "coalition": "red",
                     "range_km": 35, "location": "Kobuleti"}],
        "flights": [{"callsign": "ENFIELD", "aircraft": "FA-18C", "count": 4,
                     "role": "cas", "frequency": "305.000", "tacan": "",
                     "home_plate": "Senaki"}],
        "comms": [{"label": "STRIKE", "value": "270.800"}],
    }


def _make_base_template(n_slides: int = 2) -> str:
    """A small styled 16:9 .pptx (n blank slides), returned base64 — stands
    in for a squadron template upload."""
    import base64
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(n_slides):
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = f"COVER {i+1}"
    out = io.BytesIO(); prs.save(out)
    return base64.b64encode(out.getvalue()).decode("ascii")


def _make_cover_template() -> str:
    """A template with a 'Title'/'Sub Title' cover slide + a blank 2nd slide
    — mirrors the squadron template shape that prompted the fixes."""
    import base64
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = "Title"
    s1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1)).text_frame.text = "Sub Title"
    prs.slides.add_slide(blank)  # blank 2nd slide (no text)
    out = io.BytesIO(); prs.save(out)
    return base64.b64encode(out.getvalue()).decode("ascii")


class TestRenderWingBriefBaseTemplate:
    def test_default_render_has_ten_slides(self):
        from services.brief_renderer import render_wing_brief
        prs = Presentation(io.BytesIO(render_wing_brief(_minimal_wing_brief())))
        assert len(prs.slides._sldIdLst) == 10  # built-in cover + 9 content

    def test_base_template_keeps_its_cover_and_drops_builtin(self):
        from services.brief_renderer import render_wing_brief
        tpl = _make_base_template(2)
        prs = Presentation(io.BytesIO(
            render_wing_brief(_minimal_wing_brief(), base_template_b64=tpl)
        ))
        # 2 template slides + 9 content (built-in cover dropped) = 11.
        assert len(prs.slides._sldIdLst) == 11
        # The deck still opens and the first slide is the template's cover.
        first = prs.slides[0]
        text = " ".join(
            r.text for sh in first.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs
        )
        assert "COVER 1" in text

    def test_bad_template_falls_back_to_default(self):
        from services.brief_renderer import render_wing_brief
        prs = Presentation(io.BytesIO(
            render_wing_brief(_minimal_wing_brief(), base_template_b64="not-valid-base64!!")
        ))
        # Garbage template must not break the brief — falls back to default.
        assert len(prs.slides._sldIdLst) == 10

    def test_data_uri_prefixed_template_is_accepted(self):
        from services.brief_renderer import render_wing_brief
        tpl = "data:application/vnd.ms-powerpoint;base64," + _make_base_template(1)
        prs = Presentation(io.BytesIO(
            render_wing_brief(_minimal_wing_brief(), base_template_b64=tpl)
        ))
        # 1 template slide + 9 content = 10.
        assert len(prs.slides._sldIdLst) == 10


def _set_master_bg(prs, inner_fill_xml: str):
    """Inject a <p:bg> with the given fill into the presentation's master."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn
    m = prs.slide_masters[0].element
    csld = m.find(qn("p:cSld"))
    bg = parse_xml(
        f'<p:bg {nsdecls("p", "a")}><p:bgPr>{inner_fill_xml}<a:effectLst/></p:bgPr></p:bg>'
    )
    csld.insert(0, bg)


def _template_with_bg(inner_fill_xml: str) -> str:
    import base64
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    _set_master_bg(prs, inner_fill_xml)
    prs.slides.add_slide(prs.slide_layouts[6])
    out = io.BytesIO(); prs.save(out)
    return base64.b64encode(out.getvalue()).decode("ascii")


class TestPaletteDetection:
    def test_hex_is_dark(self):
        from services.brief_renderer import _hex_is_dark
        assert _hex_is_dark("000000") is True
        assert _hex_is_dark("1A1A1A") is True
        assert _hex_is_dark("FFFFFF") is False
        assert _hex_is_dark("F3F3F3") is False
        assert _hex_is_dark("#101010") is True

    def test_no_explicit_bg_assumes_light(self):
        from services.brief_renderer import _master_bg_is_dark
        prs = Presentation()  # default master has no explicit <p:bg>
        assert _master_bg_is_dark(prs) is False

    def test_dark_srgb_master_detected_dark(self):
        from services.brief_renderer import _master_bg_is_dark
        prs = Presentation()
        _set_master_bg(prs, '<a:solidFill><a:srgbClr val="101010"/></a:solidFill>')
        assert _master_bg_is_dark(prs) is True

    def test_light_srgb_master_detected_light(self):
        from services.brief_renderer import _master_bg_is_dark
        prs = Presentation()
        _set_master_bg(prs, '<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>')
        assert _master_bg_is_dark(prs) is False

    def test_scheme_lt1_is_light_dk1_is_dark(self):
        from services.brief_renderer import _master_bg_is_dark
        light = Presentation()
        _set_master_bg(light, '<a:solidFill><a:schemeClr val="lt1"/></a:solidFill>')
        assert _master_bg_is_dark(light) is False
        dk = Presentation()
        _set_master_bg(dk, '<a:solidFill><a:schemeClr val="dk1"/></a:solidFill>')
        assert _master_bg_is_dark(dk) is True

    def test_top_margin_shifts_content_down(self):
        """A larger content top-margin pushes the section header (and all
        content) further down so it clears the template's header band."""
        from services.brief_renderer import render_wing_brief
        tpl = _make_base_template(1)

        def topmost_text(margin):
            prs = Presentation(io.BytesIO(render_wing_brief(
                _minimal_wing_brief(), base_template_b64=tpl, top_margin_in=margin,
            )))
            slide = prs.slides[1]  # first content slide (template cover at 0)
            tops = [sh.top for sh in slide.shapes if sh.has_text_frame and sh.top is not None]
            return min(tops)

        assert topmost_text(2.0) > topmost_text(0.0)

    def test_threats_divider_respects_top_margin(self):
        """Regression: the threats column-header divider used to skip the
        template top-margin offset, so it floated up near the header band
        as a stray grey line. It must now shift with the margin like the
        rest of the content."""
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief()
        brief["threats"] = [
            {"name": f"S{i}", "type": "SAM", "coalition": "red",
             "range_km": 35, "location": "X"} for i in range(8)
        ]
        tpl = _make_base_template(1)

        def deepest_divider_top(margin):
            prs = Presentation(io.BytesIO(render_wing_brief(
                brief, base_template_b64=tpl, top_margin_in=margin)))
            for idx in range(len(prs.slides._sldIdLst)):
                s = prs.slides[idx]
                t = next((sh.text_frame.text for sh in s.shapes
                          if sh.has_text_frame and sh.text_frame.text), "")
                if "THREAT" in t.upper():
                    dividers = [sh.top for sh in s.shapes
                                if sh.shape_type == 1 and sh.height
                                and sh.height < Inches(0.1)
                                and sh.width and sh.width > Inches(11)
                                and sh.top is not None]
                    return max(dividers) if dividers else None
            return None

        t0, t2 = deepest_divider_top(0.0), deepest_divider_top(2.0)
        assert t0 is not None and t2 is not None
        assert t2 > t0 + Inches(1.0)  # shifted down by ~the margin

    def test_template_margin_shrinks_rows_per_page(self):
        """A bigger top-margin must reduce flights/threats per page so tables
        keep a bottom margin instead of overflowing — i.e. more pages."""
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief()
        brief["flights"] = [
            {"callsign": f"F{i}", "aircraft": "FA-18C", "count": 4, "role": "cas",
             "frequency": "305.000", "tacan": "", "home_plate": "X"} for i in range(20)
        ]
        tpl = _make_base_template(1)

        def friendly_pages(margin):
            prs = Presentation(io.BytesIO(render_wing_brief(
                brief, base_template_b64=tpl, top_margin_in=margin)))
            return sum(
                1 for idx in range(len(prs.slides._sldIdLst))
                if any(sh.has_text_frame and "FRIENDLY FORCES" in (sh.text_frame.text or "").upper()
                       for sh in prs.slides[idx].shapes)
            )

        assert friendly_pages(2.5) > friendly_pages(0.3)

    def test_template_cover_title_filled(self):
        """The template's literal 'Title'/'Sub Title' get replaced with the
        mission name + theatre/date/time line."""
        from services.brief_renderer import render_wing_brief
        prs = Presentation(io.BytesIO(
            render_wing_brief(_minimal_wing_brief(), base_template_b64=_make_cover_template())
        ))
        cover_texts = [" ".join(sh.text_frame.text.split())
                       for sh in prs.slides[0].shapes if sh.has_text_frame]
        joined = " ".join(cover_texts)
        assert "TEST OP" in joined            # mission name landed
        assert "Title" not in cover_texts      # literal 'Title' replaced
        assert "Sub Title" not in cover_texts  # literal 'Sub Title' replaced

    def test_template_blank_second_slide_dropped(self):
        """The template's blank 2nd slide must not survive as a blank page."""
        from services.brief_renderer import render_wing_brief
        prs = Presentation(io.BytesIO(
            render_wing_brief(_minimal_wing_brief(), base_template_b64=_make_cover_template())
        ))
        # cover + 9 content = 10 (blank dropped; would be 11 otherwise).
        assert len(prs.slides._sldIdLst) == 10
        # slide after the cover is real content, not an empty page.
        s1 = " ".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
        assert s1.strip() != ""

    def test_render_on_light_template_skips_black_rect(self):
        """On a light template the brief must NOT paint its dark rectangle —
        otherwise the master branding is hidden. Verify no full-slide dark
        rectangle is present on a content slide."""
        from services.brief_renderer import render_wing_brief
        from pptx.util import Emu
        tpl = _template_with_bg('<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>')
        prs = Presentation(io.BytesIO(
            render_wing_brief(_minimal_wing_brief(), base_template_b64=tpl)
        ))
        # A content slide (last one = Notes). No autoshape should fill the
        # whole slide (that would be our dark bg rect, which we skip).
        notes_slide = prs.slides[len(prs.slides._sldIdLst) - 1]
        full = [sh for sh in notes_slide.shapes
                if sh.width and sh.height
                and sh.width >= prs.slide_width and sh.height >= prs.slide_height]
        assert full == []


# --- Air threats analysis + slide -------------------------------------------

class TestBuildAirThreats:
    def test_aggregates_by_type_excludes_friendly_and_ground(self):
        from services.brief_builder import _build_air_threats
        groups = [
            {"category": "plane", "coalition": "red",
             "units": [{"type": "Su-27"}, {"type": "Su-27"}, {"type": "Su-27"}]},
            {"category": "plane", "coalition": "red", "units": [{"type": "Su-27"}]},  # same type, 2nd group
            {"category": "plane", "coalition": "red", "units": [{"type": "MiG-29S"}]},
            {"category": "plane", "coalition": "blue", "units": [{"type": "FA-18C_hornet"}]},  # friendly
            {"category": "vehicle", "coalition": "red", "units": [{"type": "Ural"}]},          # ground
        ]
        rows = _build_air_threats(groups)
        # Su-27 aggregated across both red groups -> 4; highest count sorts first.
        assert rows[0]["composition"].startswith("4") and rows[0]["composition"].endswith("Su-27")
        assert "fighter" in rows[0]["airframe_class"].lower()
        assert "R-27" in rows[0]["weapons"]      # capability rundown, not a .miz position
        assert rows[0]["notes"]                  # non-empty tactical note
        comps = [r["composition"] for r in rows]
        assert any(c.endswith("MiG-29S") for c in comps)
        assert all("FA-18" not in c and "Ural" not in c for c in comps)

    def test_no_enemy_air_returns_empty(self):
        from services.brief_builder import _build_air_threats
        groups = [{"category": "plane", "coalition": "blue", "units": [{"type": "FA-18C_hornet"}]}]
        assert _build_air_threats(groups) == []

    def test_unknown_type_gets_generic_profile(self):
        from services.brief_builder import _build_air_threats
        rows = _build_air_threats([
            {"category": "plane", "coalition": "red", "units": [{"type": "ZZ-99_mystery"}]},
        ])
        assert len(rows) == 1
        assert rows[0]["composition"].startswith("1")
        assert rows[0]["airframe_class"] == "Unknown type"
        assert "verify" in (rows[0]["weapons"] + rows[0]["notes"]).lower()


class TestAirThreatsSlide:
    def test_air_threats_slide_added_when_present(self):
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief()
        brief["air_threats"] = [
            {"composition": "8x Su-27", "airframe_class": "Heavy 4th-gen fighter",
             "weapons": "R-27ER/ET BVR", "notes": "Defend the R-27 early.", "coalition": "red"},
        ]
        prs = Presentation(io.BytesIO(render_wing_brief(brief)))
        assert len(prs.slides._sldIdLst) == 11  # default 10 + 1 air-threats slide
        joined = "\n".join(
            sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame
        )
        assert "AIR THREATS" in joined
        assert "Su-27" in joined and "R-27" in joined

    def test_no_air_threats_keeps_default_slide_count(self):
        from services.brief_renderer import render_wing_brief
        prs = Presentation(io.BytesIO(render_wing_brief(_minimal_wing_brief())))
        assert len(prs.slides._sldIdLst) == 10

    def test_first_threats_slide_is_titled_surface_threats(self):
        from services.brief_renderer import render_wing_brief
        prs = Presentation(io.BytesIO(render_wing_brief(_minimal_wing_brief())))
        joined = "\n".join(
            sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame
        )
        assert "SURFACE THREATS" in joined


def _player_strike_mission() -> dict:
    return {
        "overview": {"start_time": 28800, "date": "2026-07-09"},
        "airbases": [{"name": "Senaki", "lat": 42.24, "lon": 42.05}],
        "groups": [{
            "groupName": "Bengal", "category": "plane", "coalition": "blue",
            "task": "strike", "frequency": 305000000,
            "units": [{"type": "FA-18C_hornet", "name": "Bengal 1", "skill": "Player"}],
            "waypoints": [
                {"waypoint_name": "DEP", "eta_seconds": 0, "lat": 42.24, "lon": 42.05, "altitude_m": 0, "speed_ms": 150},
                {"waypoint_name": "IP ALPHA", "eta_seconds": 900, "lat": 42.4, "lon": 42.3, "altitude_m": 6000, "speed_ms": 230, "leg_distance_nm": 20},
                {"waypoint_name": "TARGET", "eta_seconds": 1500, "lat": 42.5, "lon": 42.5, "altitude_m": 5000, "speed_ms": 250, "leg_distance_nm": 12},
                {"waypoint_name": "RTB", "eta_seconds": 3000, "lat": 42.24, "lon": 42.05, "altitude_m": 0, "speed_ms": 200, "leg_distance_nm": 30},
            ],
        }],
    }


class TestWingTimelineSquadronLevel:
    """The wing brief's timeline is package/squadron-level phases, NOT one row
    per flight (per-flight detail lives on each flight brief)."""

    def test_timeline_is_package_phases_not_flight_callsigns(self):
        from services.brief_builder import _build_timeline
        groups = _player_strike_mission()["groups"]
        rows = _build_timeline(28800.0, groups, "strike")
        phases = [r["phase"] for r in rows]
        assert "Takeoff" in phases and "Push" in phases and "RTB" in phases
        assert not any("Bengal" in p for p in phases)  # not per-flight rows


class TestFlightBriefSchedule:
    """Each flight brief carries its OWN schedule (Takeoff / Push / TOT / RTB
    from its own waypoints)."""

    def test_flight_brief_has_own_schedule(self):
        from services.brief_builder import build_flight_briefs
        briefs = build_flight_briefs(
            mission_data=_player_strike_mission(), theater="Caucasus", filename="x.miz",
        )
        assert len(briefs) == 1
        phases = [r["phase"] for r in briefs[0]["timeline"]]
        assert phases[0] == "Takeoff"
        assert "TOT" in phases
        assert phases[-1] == "RTB"

    def test_flight_brief_renders_with_schedule(self):
        from services.brief_builder import build_flight_briefs
        from services.brief_renderer import render_flight_brief
        briefs = build_flight_briefs(
            mission_data=_player_strike_mission(), theater="Caucasus", filename="x.miz",
        )
        prs = Presentation(io.BytesIO(render_flight_brief(briefs[0])))
        joined = "\n".join(
            sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame
        )
        assert "SCHEDULE" in joined and "TOT" in joined


class TestScenarioFleshOut:
    """The auto-built scenario synthesises Situation / Friendly Forces /
    Adversary from the mission, with the .miz's own text leading."""

    def test_synthesises_situation_friendly_adversary(self):
        from services.brief_builder import _build_scenario
        groups = [
            {"category": "plane", "coalition": "blue", "task": "strike",
             "units": [{"type": "FA-18C_hornet", "skill": "Player"}]},
            {"category": "plane", "coalition": "red", "task": "CAP",
             "units": [{"type": "Su-27"}, {"type": "Su-27"}]},
        ]
        threats = [{"name": "SA-11"}, {"name": "SA-11"}]
        out = _build_scenario(
            {"date": "2026-07-09", "start_time": 28800, "description": "Retake the field."},
            {}, groups=groups, threats=threats, theater="Caucasus",
        )
        assert "SITUATION" in out and "FRIENDLY FORCES" in out and "ADVERSARY" in out
        assert "Retake the field." in out   # mission's own text leads
        assert "F/A-18C" in out             # friendly package
        assert "Su-27" in out               # enemy air
        assert "SA-11" in out               # surface threat

    def test_empty_mission_still_produces_sections(self):
        from services.brief_builder import _build_scenario
        out = _build_scenario({}, {}, groups=[], threats=[], theater="Caucasus")
        assert "ADVERSARY" in out
        assert "no enemy aircraft detected" in out.lower()


# ---------------------------------------------------------------------------
# POPUP ATTACK slide (v1.17.6 — auto brief grows a Popup Attack section when
# the planner has profiles defined in the Kneeboard tab; v1.17.7 adds the
# per-row mini side-profile chart).
# ---------------------------------------------------------------------------

def _popup_profile(attack_type: str = "type1", name: str = "Attack 1") -> dict:
    """Frontend's PopupAttackInput dict — full geometry so the renderer's
    mini-chart helper can compute reference points for all six types."""
    return {
        "attackType": attack_type, "name": name,
        "targetElevationFt": 100, "vipDistanceNm": 8,
        "popupAltitudeFtMsl": 8000, "popupAngleDeg": 40, "angleOffsetDeg": 25,
        "diveAngleDeg": 30,
        "releaseAltitudeFtAgl": 2000, "releaseSpeedKts": 480,
        "ingressAltitudeFtAgl": 500, "ingressSpeedKts": 480,
        "recoveryAltitudeFtAgl": 500,
    }


def _slide_text_all(pptx_bytes: bytes) -> list[str]:
    """All paragraph text from every slide, one string per slide."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    out = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    parts.append("".join(r.text for r in p.runs))
        out.append("\n".join(parts))
    return out


class TestPopupAttackSlide:
    def test_empty_profiles_no_extra_slide(self):
        """Empty popup_attacks → existing 10-slide layout unchanged."""
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief() | {"popup_attacks": []}
        prs = Presentation(io.BytesIO(render_wing_brief(brief)))
        assert len(prs.slides._sldIdLst) == 10

    def test_single_profile_adds_one_slide(self):
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief() | {"popup_attacks": [_popup_profile()]}
        prs = Presentation(io.BytesIO(render_wing_brief(brief)))
        assert len(prs.slides._sldIdLst) == 11

    def test_pagination_4_per_slide(self):
        """5 profiles → 2 popup slides (4 + 1) on top of the 10 base slides."""
        from services.brief_renderer import render_wing_brief
        profiles = [_popup_profile(name=f"P{i}") for i in range(5)]
        brief = _minimal_wing_brief() | {"popup_attacks": profiles}
        prs = Presentation(io.BytesIO(render_wing_brief(brief)))
        assert len(prs.slides._sldIdLst) == 12

    def test_slide_header_pagination_label(self):
        """Multi-page popup-attack slides label themselves '(N/M)'."""
        from services.brief_renderer import render_wing_brief
        profiles = [_popup_profile(name=f"P{i}") for i in range(5)]
        brief = _minimal_wing_brief() | {"popup_attacks": profiles}
        slides = _slide_text_all(render_wing_brief(brief))
        # Slides 11/12 (0-indexed 10/11) are the popup-attack pages.
        assert "POPUP ATTACK (1/2)" in slides[10]
        assert "POPUP ATTACK (2/2)" in slides[11]

    def test_single_page_no_pagination_label(self):
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief() | {"popup_attacks": [_popup_profile()]}
        slides = _slide_text_all(render_wing_brief(brief))
        assert "POPUP ATTACK" in slides[10]
        # Header is bare 'POPUP ATTACK' — no '(1/1)' on a single-page section.
        assert "(1/1)" not in slides[10]

    def test_profile_data_renders_as_text(self):
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief() | {"popup_attacks": [
            _popup_profile(name="LGB Run"),
        ]}
        slide_text = _slide_text_all(render_wing_brief(brief))[10]
        assert "LGB Run" in slide_text
        assert "Type 1 Popup" in slide_text
        # Parameters appear in the two-column ladder.
        assert "8,000 ft MSL" in slide_text   # popup alt
        assert "2,000 ft AGL" in slide_text   # release alt
        assert "40°" in slide_text             # popup angle
        assert "30°" in slide_text             # dive angle / offset overlap

    @pytest.mark.parametrize("attack_type", [
        "type1", "type2", "type3", "laydown", "loft", "dive",
    ])
    def test_all_attack_types_render(self, attack_type):
        """Every supported type renders without throwing — covers the mini
        side-profile helper's per-type geometry branches."""
        from services.brief_renderer import render_wing_brief
        brief = _minimal_wing_brief() | {
            "popup_attacks": [_popup_profile(attack_type=attack_type)],
        }
        out = render_wing_brief(brief)
        prs = Presentation(io.BytesIO(out))
        assert len(prs.slides._sldIdLst) == 11
        # Type label maps to a human-readable string in the slide chip.
        labels = {
            "type1": "Type 1 Popup", "type2": "Type 2 Popup",
            "type3": "Type 3 Popup", "laydown": "Lay-Down",
            "loft": "Loft (Toss)", "dive": "Straight Dive",
        }
        assert labels[attack_type] in _slide_text_all(out)[10]

    def test_degenerate_profile_renders(self):
        """Bad inputs (zero angles, missing fields) should clamp, not crash."""
        from services.brief_renderer import render_wing_brief
        weird = {
            "attackType": "type1", "name": "Weird",
            "targetElevationFt": 0, "vipDistanceNm": 0,
            "popupAltitudeFtMsl": 0, "popupAngleDeg": 0,
            "diveAngleDeg": 0, "angleOffsetDeg": 0,
            "releaseAltitudeFtAgl": 0, "releaseSpeedKts": 0,
            "ingressAltitudeFtAgl": 0, "ingressSpeedKts": 0,
        }
        brief = _minimal_wing_brief() | {"popup_attacks": [weird]}
        # Must not raise; output is still a valid pptx.
        out = render_wing_brief(brief)
        prs = Presentation(io.BytesIO(out))
        assert len(prs.slides._sldIdLst) == 11

    def test_builder_passthrough(self):
        """build_wing_brief accepts popup_attacks kwarg and the dict it
        returns contains the list verbatim."""
        from services.brief_builder import build_wing_brief
        prof = _popup_profile()
        out = build_wing_brief(
            mission_data={"overview": {}, "groups": [], "threats": [], "airbases": []},
            theater="Caucasus", filename="t.miz",
            popup_attacks=[prof],
        )
        assert out["popup_attacks"] == [prof]

    def test_builder_defaults_to_empty_list(self):
        from services.brief_builder import build_wing_brief
        out = build_wing_brief(
            mission_data={"overview": {}, "groups": [], "threats": [], "airbases": []},
            theater="Caucasus", filename="t.miz",
        )
        assert out["popup_attacks"] == []


def _minimal_flight_brief() -> dict:
    """Smallest FlightBrief dict render_flight_brief accepts. Mirrors the
    fields the dataclass requires; popup_attacks defaults to empty."""
    return {
        "mission_name": "TEST OP", "theater": "Caucasus",
        "date": "2026-05-21", "time_zulu": "0830Z",
        "callsign": "ENFIELD", "aircraft": "FA-18C", "count": 4,
        "role": "cas", "home_plate": "Senaki", "divert": "Kobuleti",
        "tasking": "Strike", "waypoints": [], "frequency": "305.000",
        "tacan": "", "icls": "",
        "fuel_joker_lbs": 4500, "fuel_bingo_lbs": 3500, "fuel_rtb_lbs": 2500,
        "notes": "", "timeline": [],
    }


class TestPopupAttackFlightSlide:
    def test_empty_profiles_no_extra_slide(self):
        """No popup_attacks → flight brief is the original 4 slides
        (cover + comms+fuel + ... no notes since notes is empty)."""
        from services.brief_renderer import render_flight_brief
        out = render_flight_brief(_minimal_flight_brief() | {"popup_attacks": []})
        prs = Presentation(io.BytesIO(out))
        # Base layout has 3 slides (cover, comms+fuel, route) — notes
        # only renders when non-empty.
        baseline = len(prs.slides._sldIdLst)
        # And adding a profile bumps by exactly 1.
        out2 = render_flight_brief(_minimal_flight_brief() | {"popup_attacks": [_popup_profile()]})
        prs2 = Presentation(io.BytesIO(out2))
        assert len(prs2.slides._sldIdLst) == baseline + 1

    def test_5_profiles_paginate(self):
        from services.brief_renderer import render_flight_brief
        profiles = [_popup_profile(name=f"P{i}") for i in range(5)]
        baseline = len(Presentation(io.BytesIO(
            render_flight_brief(_minimal_flight_brief()))).slides._sldIdLst)
        prs = Presentation(io.BytesIO(
            render_flight_brief(_minimal_flight_brief() | {"popup_attacks": profiles})))
        assert len(prs.slides._sldIdLst) == baseline + 2  # 4+1 across 2 slides

    @pytest.mark.parametrize("attack_type", [
        "type1", "type2", "type3", "laydown", "loft", "dive",
    ])
    def test_per_flight_all_types_render(self, attack_type):
        from services.brief_renderer import render_flight_brief
        out = render_flight_brief(_minimal_flight_brief() | {
            "popup_attacks": [_popup_profile(attack_type=attack_type)],
        })
        # Must produce a valid pptx + include the type label.
        assert b"PK" == out[:2]  # zip header — pptx is a zip
        prs = Presentation(io.BytesIO(out))
        last_slide_text = []
        for shape in prs.slides[len(prs.slides) - 1].shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    last_slide_text.append("".join(r.text for r in p.runs))
        labels = {
            "type1": "Type 1 Popup", "type2": "Type 2 Popup",
            "type3": "Type 3 Popup", "laydown": "Lay-Down",
            "loft": "Loft (Toss)", "dive": "Straight Dive",
        }
        assert labels[attack_type] in "\n".join(last_slide_text)

    def test_flight_builder_passthrough(self):
        from services.brief_builder import build_flight_briefs
        # Need a minimal mission with a player flight.
        mission = {
            "overview": {"start_time": 0},
            "groups": [{
                "groupName": "ENFIELD", "task": "cas", "category": "plane",
                "coalition": "blue",
                "units": [{"name": "ENFIELD11", "type": "FA-18C", "skill": "Player"}],
                "waypoints": [],
            }],
            "airbases": [],
        }
        profiles = [_popup_profile()]
        out = build_flight_briefs(
            mission_data=mission, theater="Caucasus", filename="x.miz",
            popup_attacks=profiles,
        )
        assert len(out) == 1
        assert out[0]["popup_attacks"] == profiles
